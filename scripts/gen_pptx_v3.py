#!/usr/bin/env python3
"""桜花賞徹底解析 v3
Pattern1(円グラフ/ランキング) + Pattern2(大数字/番号バッジ)
+ ギーニョ / ゴトー キャラクター
12スライド構成（前半ブロック）
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

W_IN, H_IN = 13.33, 7.5

# ══ カラーパレット ══════════════════════════════
C_PINK    = RGBColor(232, 19,110)   # ビビッドピンク（主役）
C_PINK_DK = RGBColor(180,  8, 75)   # ダークピンク
C_CREAM   = RGBColor(255,252,248)   # クリーム背景
C_WARM    = RGBColor(255,245,235)   # ウォームクリーム
C_WHITE   = RGBColor(255,255,255)
C_GOLD    = RGBColor(245,183, 49)   # ゴールド
C_GOLD_DK = RGBColor(180, 83,  9)
C_GREEN   = RGBColor( 22,163, 74)   # 緑（複勝）
C_GREEN_LT= RGBColor(220,252,231)
C_TEAL    = RGBColor(  8,145,178)   # ティール（Pattern2）
C_RED     = RGBColor(220, 38, 38)   # 赤（警告）
C_ORANGE  = RGBColor(234, 88, 12)   # オレンジ
C_PURPLE  = RGBColor(109, 40,217)   # 穴馬
C_DARK    = RGBColor( 15, 10, 12)   # 黒テキスト
C_GRAY    = RGBColor(107, 80, 90)   # グレーテキスト
C_DECO    = RGBColor(245,200,170)   # 装飾シェブロン

FONT_S = "游明朝"
FONT_G = "游ゴシック"

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_BRAIN = os.path.join(BASE, 'brain_cover.jpg')
IMG_GINYO = os.path.join(BASE, 'ginyo.jpg')
IMG_GOTO  = os.path.join(BASE, 'goto.jpg')

# ══ データ ═════════════════════════════════════
with open(os.path.join(BASE,'docs/data/race-notes/2026-04-12-hanshin-11r.json'),encoding='utf-8') as f:
    DATA = json.load(f)
RACE_ANA  = DATA.get('raceAnalysis','')
BLOOD_ANA = DATA.get('bloodAnalysis','')
HORSES    = DATA['horses']

GP = {'3S':12,'2S':11,'1S':10,'3A':9,'2A':8,'1A':7,'3B':6,'2B':5,'1B':4,
      '3C':3,'2C':2,'1C':1,'D':0,'S':11,'A':8,'B':5,'C':2}
GM = {'zi':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'ext':{'S':['A','S'],'A':['B','A'],'B':['C','B']},
      'kinso':{'S':['S','S'],'A':['A','A'],'B':['B','A'],'C':['D','C'],'D':['D','D']},
      'pace':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'body':{'S':['A','S'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'course':{'S':['A','S'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'style':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']}}

def zi_g(v):
    if not v: return 'B'
    return 'S' if v>=130 else 'A' if v>=115 else 'B' if v>=100 else 'C' if v>=85 else 'D'
def ext_g(f):
    if not f or f=='在厩調整': return 'B'
    if '天栄' in f or 'しがらき' in f: return 'S'
    if 'チャンピオン' in f or '山元' in f or 'キャニオン' in f: return 'A'
    return 'B'
def splt(g,k):
    if not g or g=='—': return [None,None]
    if k in GM and g in GM[k]: return list(GM[k][g])
    return [g,g] if g else [None,None]
def gp(g): return GP.get(g) if g else None

def calc(h):
    zig=zi_g(h.get('sabcZI')); exg=ext_g(h.get('extFacility'))
    zt,zf=splt(zig,'zi'); et,ef=splt(exg,'ext')
    kt,kf=splt(h.get('kinsoGrade'),'kinso'); pt,pf=splt(h.get('paceGrade'),'pace')
    bdt,bdf=splt(h.get('bodyGrade'),'body'); ct,cf=splt(h.get('courseGrade'),'course')
    syt,syf=splt(h.get('styleGrade'),'style')
    srt=h.get('sireGradeTan'); srf=h.get('sireGradeFuku')
    dst=h.get('damSireGradeTan'); dsf=h.get('damSireGradeFuku')
    sr2t=h.get('sireGrade2Tan','B'); sr2f=h.get('sireGrade2Fuku','B')
    ds2t=h.get('damSireGrade2Tan','B'); ds2f=h.get('damSireGrade2Fuku','B')
    blt=h.get('bloodGradeTan') if h.get('bloodGradeTan') and h.get('bloodGradeTan')!='—' else None
    blf=h.get('bloodGradeFuku') if h.get('bloodGradeFuku') and h.get('bloodGradeFuku')!='—' else None
    stt=h.get('stableGradeTan'); stf=h.get('stableGradeFuku')
    brt=h.get('breederGradeTan'); brf=h.get('breederGradeFuku')
    jkt=h.get('jockeyGradeTan'); jkf=h.get('jockeyGradeFuku')
    gtt=h.get('gateGradeTan'); gtf=h.get('gateGradeFuku')
    spt=h.get('specialGradeTan') or h.get('specialGrade')
    spf=h.get('specialGradeFuku') or h.get('specialGrade')
    agt=h.get('ageGradeTan'); agf=h.get('ageGradeFuku')
    rtt=h.get('rotGradeTan'); rtf=h.get('rotGradeFuku')
    s3t=h.get('sf3GradeTan'); s3f=h.get('sf3GradeFuku')
    cdt=h.get('conditionGradeTan'); cdf=h.get('conditionGradeFuku')
    smt=h.get('somaGradeTan'); smf=h.get('somaGradeFuku')
    tgs=[zt,et,srt,dst,sr2t,ds2t,blt,kt,pt,stt,bdt,ct,syt,brt,jkt,gtt,spt,agt,rtt,s3t,cdt,smt]
    fgs=[zf,ef,srf,dsf,sr2f,ds2f,blf,kf,pf,stf,bdf,cf,syf,brf,jkf,gtf,spf,agf,rtf,s3f,cdf,smf]
    ts=tf=tm=fm=0
    for g in tgs:
        p=gp(g)
        if p is not None: ts+=p; tm+=12
    for g in fgs:
        p=gp(g)
        if p is not None: tf+=p; fm+=12
    sf3=(h.get('sf3AnaScore') or 0); rab=(h.get('raceAnaBonus') or 0)
    bp=((gp(blt or 'B') or 5)+(gp(blf or 'B') or 5)-10)/24
    ana=sf3+rab+bp
    sougou=(ts+tf)/(tm+fm) if (tm+fm)>0 else 0
    return {'tanScore':ts,'tanMax':tm,'fukuScore':tf,'fukuMax':fm,'anaScore':ana,'sougouScore':sougou}

horse_data=[]
for name,h in HORSES.items():
    sc=calc(h)
    tp=sc['tanScore']/sc['tanMax'] if sc['tanMax']>0 else 0
    fp=sc['fukuScore']/sc['fukuMax'] if sc['fukuMax']>0 else 0
    horse_data.append({**sc,'name':name,'h':h,'tanPct':tp,'fukuPct':fp,'odds':h.get('expectedOdds',999)})

by_tan    = sorted(horse_data,key=lambda x:-x['tanPct'])
by_fuku   = sorted(horse_data,key=lambda x:-x['fukuPct'])
by_ana    = sorted(horse_data,key=lambda x:-x['anaScore'])
by_sougou = sorted(horse_data,key=lambda x:-x['sougouScore'])
for i,d in enumerate(by_tan):    d['tanRank']=i+1
for i,d in enumerate(by_fuku):   d['fukuRank']=i+1
for i,d in enumerate(by_ana):    d['anaRank']=i+1
for i,d in enumerate(by_sougou): d['sougouRank']=i+1

# ══ PPTX ヘルパー ══════════════════════════════
def new_prs():
    p=Presentation(); p.slide_width=Inches(W_IN); p.slide_height=Inches(H_IN); return p

def add_slide(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl,c):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def rect(sl,x,y,w,h,fill=None,line=None,lw=1.0,radius=False):
    s=sl.shapes.add_shape(5 if radius else 1,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    if line: s.line.color.rgb=line; s.line.width=Pt(lw)
    else: s.line.fill.background()
    return s

def txt(sl,text,x,y,w,h,size=14,bold=False,color=None,
        align=PP_ALIGN.LEFT,italic=False,font=None,wrap=True):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    if color: r.font.color.rgb=color
    fn=font or FONT_G; r.font.name=fn
    rPr=r._r.get_or_add_rPr()
    ea=rPr.find(qn('a:ea'))
    if ea is None: ea=etree.SubElement(rPr,qn('a:ea'))
    ea.set('typeface',fn)
    return tb

def pic(sl,path,x,y,w,h):
    if os.path.exists(path):
        sl.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))

def chevrons(sl):
    """Pattern1風 コーナー装飾"""
    txt(sl,'»»',0.15,0.08,1.6,0.65,size=30,bold=True,color=C_DECO,font=FONT_G)
    txt(sl,'»»',11.4,6.72,1.7,0.6,size=30,bold=True,color=C_DECO,align=PP_ALIGN.RIGHT,font=FONT_G)

def footer(sl,page,total=12,paid=False):
    rect(sl,0,7.1,13.33,0.4,fill=C_PINK_DK)
    tag='🔒 メンバーシップ限定' if paid else '🆓 無料公開'
    tag_c=RGBColor(255,230,150) if paid else C_GOLD
    txt(sl,'UG競馬 競馬予想チャンネル  ／  桜花賞徹底解析 2026',
        0.4,7.13,9.5,0.3,size=10,color=RGBColor(255,210,230))
    txt(sl,f'{tag}  {page}/{total}',9.5,7.13,3.5,0.3,size=10,bold=True,
        color=tag_c,align=PP_ALIGN.RIGHT)

def hbar(sl,x,y,bw,pct,fc,tc=None,h=0.16):
    t=tc or RGBColor(230,215,222)
    rect(sl,x,y,bw,h,fill=t)
    if pct>0: rect(sl,x,y,bw*min(pct,1.0),h,fill=fc)

# ══ Slide 01: 表紙 ══════════════════════════════
def s01(prs):
    sl=add_slide(prs)
    bg(sl,C_PINK)
    pic(sl,IMG_BRAIN,5.6,-0.1,8.0,7.7)         # 脳画像右半分
    rect(sl,0,0,6.2,7.5,fill=C_PINK)           # 左側オーバーレイ
    rect(sl,0,0,13.33,0.07,fill=C_GOLD)        # 上部ゴールドライン
    rect(sl,0,7.43,13.33,0.07,fill=C_GOLD)     # 下部ゴールドライン

    # G1 バッジ
    rect(sl,0.45,0.55,1.4,0.55,fill=C_GOLD,radius=True)
    txt(sl,'G 1',0.45,0.55,1.4,0.55,size=20,bold=True,color=C_DARK,align=PP_ALIGN.CENTER,font=FONT_S)

    # メインタイトル
    txt(sl,'桜花賞',0.4,1.22,5.7,1.85,size=80,bold=True,color=C_WHITE,font=FONT_S)

    # サブタイトル
    rect(sl,0.4,3.18,5.7,0.07,fill=C_GOLD)
    txt(sl,'徹 底 解 析',0.4,3.35,5.7,0.88,size=36,bold=True,color=C_GOLD,font=FONT_S)

    # 日時
    txt(sl,'2026年4月12日（日）  阪神競馬場  芝1600m外  19頭立て',
        0.4,4.38,5.7,0.48,size=12,color=RGBColor(255,220,235))

    rect(sl,0.4,4.98,5.7,0.05,fill=RGBColor(255,180,210))
    txt(sl,'ギーニョ重賞データ × 血統解析 × 特別ファクター',
        0.4,5.1,5.7,0.48,size=13,italic=True,color=RGBColor(255,240,248))

    # ゴトー キャラクター（右端）
    pic(sl,IMG_GOTO,9.8,3.2,3.2,2.4)

    # ロゴ帯
    rect(sl,0,6.88,6.4,0.55,fill=C_WHITE)
    txt(sl,'UG競馬 競馬予想チャンネル',0.5,6.93,5.8,0.44,size=13,bold=True,color=C_PINK_DK)

# ══ Slide 02-06: 大数字スライド ══════════════════
def big_num(prs, badge_num, badge_text, badge_col,
            number, num_col, main_label, sub_text,
            note=None, char='ginyo', page=2):
    sl=add_slide(prs)
    bg(sl,C_CREAM)
    chevrons(sl)

    # 番号バッジ（Pattern2 スタイル）
    bw=len(badge_text)*0.15+1.3
    rect(sl,0.42,0.35,bw,0.55,fill=badge_col,radius=True)
    txt(sl,f'{badge_num}  {badge_text}',0.5,0.35,bw-0.1,0.55,
        size=14,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)

    # 巨大数字（Pattern2 主役）
    txt(sl,number,0,1.15,13.33,3.85,size=110,bold=True,
        color=num_col,align=PP_ALIGN.CENTER,font=FONT_S)

    # 区切り線
    rect(sl,1.8,5.08,9.73,0.06,fill=badge_col)

    # メインラベル
    txt(sl,main_label,0,5.2,13.33,0.72,size=22,bold=True,
        color=C_DARK,align=PP_ALIGN.CENTER,font=FONT_S)

    # サブテキスト
    txt(sl,sub_text,0.5,6.02,12.33,0.55,size=14,italic=True,
        color=C_GRAY,align=PP_ALIGN.CENTER)

    # 注釈ボックス（オプション）
    if note:
        rect(sl,3.8,6.6,5.73,0.44,fill=badge_col,radius=True)
        txt(sl,note,3.9,6.6,5.5,0.44,size=12,bold=True,
            color=C_WHITE,align=PP_ALIGN.CENTER)

    # キャラクター
    if char=='ginyo': pic(sl,IMG_GINYO,10.05,4.3,3.1,2.32)
    elif char=='goto': pic(sl,IMG_GOTO,10.2,4.2,2.9,2.18)

    footer(sl,page)

def s02(prs): big_num(prs,'①','ギーニョデータ解析',C_GOLD_DK,
    '70 %',C_GOLD_DK,'2番人気の複勝率',
    '― 1番人気（60%）を上回る「最強人気ゾーン」―',
    note='2番人気 = 過去10年 最も期待値が高い人気',char='ginyo',page=2)

def s03(prs): big_num(prs,'②','ギーニョデータ解析',C_GREEN,
    '100 %',C_GREEN,'G1前走1着の複勝率',
    '― 阪神JFからの直行ローテ  過去10年 全馬複勝圏 ―',
    note='前代未聞のデータ  →  スターアニスに直結',char='ginyo',page=3)

def s04(prs): big_num(prs,'③','ギーニョデータ解析',C_PINK,
    '80 %',C_PINK,'上がり3F  1位馬の複勝率',
    '― 直線で一番切れる馬が主役になるレース ―',
    note='上がり3F最速馬 →  勝率40% / 複勝率80%',char='goto',page=4)

def s05(prs): big_num(prs,'④','ギーニョデータ解析',C_RED,
    '0 %',C_RED,'10番人気以下の複勝率',
    '― 人気薄は完全消し。上位人気に絞り込む ―',
    note='穴馬を買う必要はゼロ  →  上位8番人気内に集中',char='ginyo',page=5)

def s06(prs): big_num(prs,'⑤','ギーニョデータ解析',C_ORANGE,
    '5 %',C_ORANGE,'G2前走1着の複勝率',
    '― チューリップ賞前走は見かけ上強そうでも危険ゾーン ―',
    note='G2前走1着 = 過去10年 最大の罠データ',char='goto',page=6)

# ══ Slide 07: 血統解析 ══════════════════════════
def s07(prs):
    sl=add_slide(prs)
    bg(sl,C_WARM)
    chevrons(sl)

    # タイトルバッジ（Pattern1 スタイル）
    rect(sl,1.5,0.22,10.33,0.72,fill=C_PINK,radius=True)
    txt(sl,'血 統 解 析',1.5,0.22,10.33,0.72,size=24,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,'血 統 解 析',1.55,0.22,10.23,0.72,size=24,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,'ギーニョ血統データ  桜花賞 父系・母父系 複勝率ランキング',
        0.5,1.05,12.33,0.4,size=12,italic=True,color=C_GRAY,align=PP_ALIGN.CENTER)

    # 左列: 父系
    rect(sl,0.4,1.55,5.9,0.52,fill=C_GOLD_DK)
    txt(sl,'父　系  ランキング',0.4,1.55,5.9,0.52,size=15,bold=True,
        color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)

    sires=[
        ('キタサンブラック','66.7%',C_GOLD,'3頭','ブラックチャリス'),
        ('エピファネイア',  '28.6%',C_PINK,'7頭','アランカール'),
        ('ロードカナロア',  '22.2%',C_GRAY,'多数','複は〇 単は✕'),
    ]
    for i,(nm,pct,col,cnt,note) in enumerate(sires):
        ry=2.15+i*1.48
        pv=float(pct.replace('%',''))
        bg_c=RGBColor(255,250,240) if i%2==0 else C_WHITE
        rect(sl,0.4,ry,5.9,1.35,fill=bg_c,radius=True)
        rect(sl,0.4,ry,0.07,1.35,fill=col)
        rnk=['🥇','🥈','🥉'][i]
        txt(sl,f'{rnk}  {nm}',0.6,ry+0.1,3.5,0.48,size=14,bold=(i==0),color=C_DARK,font=FONT_S)
        txt(sl,f'複勝率  {pct}  ({cnt})',0.6,ry+0.62,3.5,0.35,size=10,color=C_GRAY)
        hbar(sl,4.28,ry+0.2,1.75,pv/100,col,h=0.25)
        txt(sl,note,0.6,ry+0.98,5.1,0.32,size=9,italic=True,color=col)

    # 右列: 母父系
    rect(sl,7.0,1.55,5.9,0.52,fill=C_TEAL)
    txt(sl,'母父系  ランキング',7.0,1.55,5.9,0.52,size=15,bold=True,
        color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)

    dams=[
        ('ハーツクライ',    '50.0%',C_GOLD,'2頭','フェスティバルヒル'),
        ('クロフネ',        '42.9%',C_TEAL,'7頭',''),
        ('ディープインパクト','0.0%',C_RED, '多数','ルールザウェイヴ 消し'),
    ]
    for i,(nm,pct,col,cnt,note) in enumerate(dams):
        ry=2.15+i*1.48
        pv=float(pct.replace('%',''))
        bg_c=RGBColor(240,252,255) if i%2==0 else C_WHITE
        rect(sl,7.0,ry,5.9,1.35,fill=bg_c,radius=True)
        rect(sl,7.0,ry,0.07,1.35,fill=col)
        rnk=['🥇','🥈','🥉'][i]
        txt(sl,f'{rnk}  {nm}',7.2,ry+0.1,3.4,0.48,size=14,bold=(i==0),color=C_DARK,font=FONT_S)
        txt(sl,f'複勝率  {pct}  ({cnt})',7.2,ry+0.62,3.4,0.35,size=10,color=C_GRAY)
        hbar(sl,10.75,ry+0.2,1.88,pv/100,col,h=0.25)
        if note:
            txt(sl,note,7.2,ry+0.98,5.5,0.32,size=9,italic=True,color=col)

    # 血統総括 + ギーニョ
    rect(sl,0.4,6.58,8.5,0.43,fill=C_WHITE,radius=True)
    rect(sl,0.4,6.58,0.07,0.43,fill=C_GOLD)
    txt(sl,'血統総括: ブラックチャリス（父66.7%）とフェスティバルヒル（母父50%）が血統最上位。ただしデータ相克あり。',
        0.6,6.62,8.2,0.34,size=9,color=C_DARK,wrap=True)
    pic(sl,IMG_GINYO,9.1,5.4,3.9,2.92)  # ギーニョ右下
    footer(sl,7)

# ══ Slide 08: 特別解析（3列アイコンカード） ══════
def s08(prs):
    sl=add_slide(prs)
    bg(sl,C_CREAM)
    chevrons(sl)

    # タイトル
    rect(sl,1.5,0.22,10.33,0.72,fill=C_PINK,radius=True)
    txt(sl,'特 別 解 析',1.55,0.22,10.23,0.72,size=24,bold=True,
        color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,'G2前走の罠  /  G1前走の正解  /  穴馬ファクター該当馬',
        0.5,1.05,12.33,0.4,size=12,italic=True,color=C_GRAY,align=PP_ALIGN.CENTER)

    # 3列カード（Pattern2 ティールカード スタイル）
    cards=[
        (C_PINK_DK,'🚫','G2前走の罠',
         'チューリップ賞などG2前走1着は複勝率わずか5%。\n実力があっても大幅割引が正解。',
         [n for n,h in HORSES.items() if '賞' in (h.get('prevRaceName','')) and
          (h.get('prevFinish','').startswith('1') or '1着' in h.get('prevFinish',''))
          and h.get('prevRaceName','') in ['チューリップ賞','フィリーズレビュー']][:3]),
        (C_GREEN,'✅','G1前走の正解',
         'G1前走1着は複勝率100%。\n阪神JF直行ローテは最強条件。',
         [n for n,h in HORSES.items() if 'G1' in (h.get('prevRaceName','')) and
          ('1着' in (h.get('prevFinish','')) or h.get('prevFinish','')=='1着')][:3]),
        (C_GOLD_DK,'🎯','穴馬ファクター',
         'レース統計から導かれた穴馬候補。\n血統×ローテが揃った一発期待馬。',
         [n for n,h in HORSES.items() if (h.get('raceAnaBonus') or 0)>=1.0][:4]),
    ]
    col_x=[0.4,4.62,8.84]; col_w=3.85
    for ci,(col,icon,title,desc,horses) in enumerate(cards):
        cx=col_x[ci]
        # カード背景
        rect(sl,cx,1.58,col_w,4.9,fill=col,radius=True)
        # アイコン円
        txt(sl,icon,cx+col_w/2-0.45,1.72,0.9,0.82,size=32,align=PP_ALIGN.CENTER)
        # タイトル
        txt(sl,title,cx+0.15,2.6,col_w-0.3,0.55,size=16,bold=True,
            color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)
        # 説明
        txt(sl,desc,cx+0.15,3.25,col_w-0.3,0.92,size=10,
            color=RGBColor(255,240,245) if col==C_PINK_DK else C_WHITE,
            align=PP_ALIGN.CENTER,wrap=True)
        # 馬名リスト
        for hi,hn in enumerate(horses):
            hy=4.32+hi*0.52
            rect(sl,cx+0.18,hy,col_w-0.36,0.44,fill=C_WHITE,radius=True)
            txt(sl,hn,cx+0.28,hy+0.05,col_w-0.56,0.36,size=11,bold=True,color=col)

    # ゴトー
    pic(sl,IMG_GOTO,0.3,4.8,2.1,1.58)
    footer(sl,8)

# ══ Slide 09: 総合スコア1位 ══════════════════════
def s09(prs):
    sl=add_slide(prs)
    bg(sl,C_CREAM)
    top=by_sougou[0]; h=top['h']
    chevrons(sl)

    # ヘッダー
    rect(sl,0,0,13.33,1.38,fill=C_PINK)
    txt(sl,'総合スコア 1 位',0.5,0.1,9,0.75,size=30,bold=True,color=C_WHITE,font=FONT_S)
    txt(sl,f'単勝×複勝 バランス総合評価  /  {top["name"]}  {top["odds"]}倍',
        0.5,0.88,12,0.44,size=12,italic=True,color=RGBColor(255,220,235))

    # 左: スコアビジュアル
    rect(sl,0.38,1.55,5.6,5.45,fill=C_WHITE,radius=True)
    rect(sl,0.38,1.55,5.6,0.55,fill=C_PINK)
    txt(sl,f'総合  #{top["sougouRank"]}',0.55,1.6,5.25,0.45,size=13,bold=True,color=C_WHITE)
    txt(sl,top['name'],0.45,2.2,5.4,0.95,size=28,bold=True,
        color=C_PINK,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,f'想定オッズ  {top["odds"]}倍',0.45,3.22,5.4,0.38,size=13,
        color=C_GRAY,align=PP_ALIGN.CENTER)

    for i,(lbl,pct,col,rnk) in enumerate([
        ('総合',top['sougouScore'],C_PINK,top['sougouRank']),
        ('単勝',top['tanPct'],C_GOLD_DK,top['tanRank']),
        ('複勝',top['fukuPct'],C_GREEN,top['fukuRank'])]):
        by=3.75+i*0.68
        txt(sl,lbl,0.55,by,0.7,0.46,size=10,bold=True,color=col)
        hbar(sl,1.32,by+0.12,3.55,pct,col,h=0.24)
        txt(sl,f'{pct*100:.1f}%  #{rnk}',4.92,by,0.92,0.46,size=9,bold=True,
            color=col,align=PP_ALIGN.RIGHT)

    sp=h.get('specialNote','')
    if sp:
        rect(sl,0.55,5.95,5.28,0.88,fill=RGBColor(255,245,252))
        rect(sl,0.55,5.95,0.07,0.88,fill=C_GOLD)
        txt(sl,'⭐ '+sp,0.72,5.99,5.0,0.78,size=8,italic=True,color=C_DARK,wrap=True)

    # 右: ファクター + ゴトー
    rect(sl,6.28,1.55,6.7,5.45,fill=C_WHITE,radius=True)
    rect(sl,6.28,1.55,6.7,0.55,fill=C_PINK_DK)
    txt(sl,'ファクター詳細',6.45,1.6,6.35,0.45,size=13,bold=True,color=C_WHITE)

    factors=[
        ('能力 ZI',str(h.get('sabcZI','?'))),
        ('近走グレード',str(h.get('kinsoGrade','—'))),
        ('コース適性',str(h.get('courseGrade','—'))),
        ('ローテ評価',str(h.get('rotGradeTan','—'))),
        ('厩舎',str(h.get('stableGradeTan','—'))),
        ('調教',str(h.get('conditionGradeTan','—'))),
        ('前走',str(h.get('prevRaceName','—'))),
        ('前走結果',str(h.get('prevFinish','—'))),
    ]
    for i,(lbl,val) in enumerate(factors):
        fy=2.22+i*0.6
        bc=RGBColor(253,245,250) if i%2==0 else C_WHITE
        rect(sl,6.28,fy,6.7,0.55,fill=bc)
        txt(sl,lbl,6.42,fy+0.1,2.6,0.36,size=10,color=C_GRAY)
        good=any(v in val for v in ['S','2S','3S','A','2A','3A','1着','最上位'])
        vc=C_PINK if good else C_DARK
        txt(sl,val,9.12,fy+0.1,3.72,0.36,size=11,bold=good,color=vc,align=PP_ALIGN.RIGHT)

    pic(sl,IMG_GOTO,10.0,4.62,3.1,2.32)
    footer(sl,9,paid=True)

# ══ Slide 10: 単勝スコア1位 ═════════════════════
def s10(prs):
    sl=add_slide(prs)
    bg(sl,C_CREAM)
    top=by_tan[0]; h=top['h']
    chevrons(sl)

    # ゴールドヘッダー
    rect(sl,0,0,13.33,1.38,fill=C_GOLD_DK)
    txt(sl,'単勝スコア 1 位',0.5,0.1,9,0.75,size=30,bold=True,color=C_WHITE,font=FONT_S)
    txt(sl,f'勝ち馬スコア最上位  /  {top["name"]}  {top["odds"]}倍',
        0.5,0.88,12,0.44,size=12,italic=True,color=RGBColor(255,245,210))

    # 大きなスコア表示（Pattern2 ミックス）
    rect(sl,0.38,1.55,5.6,2.55,fill=C_WHITE,radius=True)
    txt(sl,top['name'],0.45,1.68,5.4,0.92,size=26,bold=True,
        color=C_GOLD_DK,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,f'{top["tanPct"]*100:.1f} %',0.45,2.62,5.4,0.88,size=44,bold=True,
        color=C_GOLD_DK,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,'単勝スコア',0.45,3.55,5.4,0.38,size=12,color=C_GRAY,align=PP_ALIGN.CENTER)
    hbar(sl,0.65,4.05,5.12,top['tanPct'],C_GOLD_DK,h=0.3)

    # なぜ勝てるか 3要素
    rect(sl,0.38,4.48,5.6,2.42,fill=C_WHITE,radius=True)
    rect(sl,0.38,4.48,5.6,0.42,fill=C_GOLD_DK)
    txt(sl,'単勝1位の理由',0.55,4.52,5.25,0.35,size=12,bold=True,color=C_WHITE)
    reasons=[]
    if h.get('sabcZI',0)>=130: reasons.append(f'ZI {h["sabcZI"]}  全馬最高値')
    if (h.get('rotGradeTan') or '') in ['S','2S','3S']:
        reasons.append(f'ローテ {h.get("rotGradeTan")}  理想的ローテ')
    if (h.get('conditionGradeTan') or '') in ['S','2S','3S']:
        reasons.append(f'調教 {h.get("conditionGradeTan")}  万全仕上がり')
    if (h.get('courseGrade') or '') in ['S','A']:
        reasons.append(f'コース適性 {h.get("courseGrade")}  阪神マイル向き')
    for ri,r in enumerate(reasons[:4]):
        rect(sl,0.55,4.98+ri*0.46,5.24,0.4,fill=RGBColor(255,250,230),radius=True)
        txt(sl,f'✓ {r}',0.7,5.02+ri*0.46,4.95,0.32,size=11,color=C_GOLD_DK)

    # 右: 騎手コメント + 特注
    rect(sl,6.28,1.55,6.7,5.45,fill=C_WHITE,radius=True)
    rect(sl,6.28,1.55,6.7,0.52,fill=C_GOLD_DK)
    txt(sl,'騎手コメント（前走）',6.45,1.6,6.35,0.42,size=12,bold=True,color=C_WHITE)
    pjc=h.get('prevJockeyComment','')
    txt(sl,pjc[:230]+'…' if len(pjc)>230 else pjc,
        6.45,2.18,6.42,2.75,size=10,italic=True,color=C_DARK,wrap=True)
    rect(sl,6.28,5.0,6.7,0.07,fill=RGBColor(230,190,90))
    sp=h.get('specialNote','')
    if sp:
        rect(sl,6.28,5.12,6.7,1.68,fill=RGBColor(255,251,232))
        rect(sl,6.28,5.12,0.07,1.68,fill=C_GOLD)
        txt(sl,'⭐ '+sp,6.45,5.16,6.42,1.58,size=9,italic=True,color=C_DARK,wrap=True)

    pic(sl,IMG_GINYO,0.1,5.7,3.1,2.32)  # ギーニョ左下
    footer(sl,10,paid=True)

# ══ Slide 11: 穴馬スコア1位 ═════════════════════
def s11(prs):
    sl=add_slide(prs)
    bg(sl,C_CREAM)
    top=by_ana[0]; h=top['h']
    chevrons(sl)

    # パープルヘッダー
    rect(sl,0,0,13.33,1.38,fill=C_PURPLE)
    txt(sl,'穴馬スコア 1 位',0.5,0.1,9,0.75,size=30,bold=True,color=C_WHITE,font=FONT_S)
    txt(sl,f'穴馬スコア最上位  /  {top["name"]}  {top["odds"]}倍  （単{top["tanRank"]}位）',
        0.5,0.88,12,0.44,size=12,italic=True,color=RGBColor(220,200,255))

    # 左: データの相克カード
    rect(sl,0.38,1.55,5.6,5.45,fill=C_WHITE,radius=True)
    txt(sl,top['name'],0.45,1.65,5.4,0.88,size=26,bold=True,
        color=C_PURPLE,align=PP_ALIGN.CENTER,font=FONT_S)
    txt(sl,f'{top["odds"]}倍  /  穴スコア {top["anaScore"]:.2f} pt',
        0.45,2.6,5.4,0.38,size=12,color=C_GRAY,align=PP_ALIGN.CENTER)
    hbar(sl,0.65,3.08,5.12,min(top['anaScore']/3.0,1.0),C_PURPLE,h=0.28)

    # 血統最高 vs レースデータ
    rect(sl,0.45,3.52,5.38,0.06,fill=RGBColor(200,180,255))
    # 血統カード
    rect(sl,0.45,3.65,5.38,1.3,fill=RGBColor(250,245,255),radius=True)
    rect(sl,0.45,3.65,0.07,1.3,fill=C_GOLD)
    txt(sl,'🧬 血統評価  最高',0.62,3.7,5.0,0.4,size=12,bold=True,color=C_GOLD_DK)
    txt(sl,'父キタサンブラック →  複勝率 66.7%\n全産駒複勝圏の驚異データ',
        0.62,4.14,5.0,0.72,size=10,color=C_DARK,wrap=True)
    # レースデータカード
    rect(sl,0.45,5.05,5.38,1.22,fill=RGBColor(255,242,242),radius=True)
    rect(sl,0.45,5.05,0.07,1.22,fill=C_RED)
    txt(sl,'📊 レースデータ  消し根拠あり',0.62,5.1,5.0,0.4,size=12,bold=True,color=C_RED)
    ran=h.get('raceAnaNote',''); sp=h.get('specialNote','')
    note_text=f'穴ボーナス: {ran}' if ran else '前走クラス等でデータ消し根拠が存在'
    txt(sl,note_text,0.62,5.54,5.0,0.62,size=10,color=C_DARK,wrap=True)

    # 右: 結論 + コメント
    rect(sl,6.28,1.55,6.7,5.45,fill=C_WHITE,radius=True)
    rect(sl,6.28,1.55,6.7,0.52,fill=C_PURPLE)
    txt(sl,'⚡ データの相克  解説',6.45,1.6,6.35,0.42,size=12,bold=True,color=C_WHITE)
    txt(sl,'血統は「最高評価」、レースデータは「消し根拠あり」という相矛盾するデータが共存。'
        'これが「データの相克」。単勝ではなく複勝での少額バックが合理的判断。',
        6.45,2.18,6.42,1.45,size=10,color=C_DARK,wrap=True)
    rect(sl,6.45,3.72,6.25,0.06,fill=RGBColor(180,150,255))
    bn=h.get('bloodNote','')
    txt(sl,'🧬 血統コメント',6.45,3.84,6.35,0.35,size=11,bold=True,color=C_PURPLE)
    txt(sl,bn[:160]+'…' if len(bn)>160 else bn,
        6.45,4.22,6.42,1.38,size=9,italic=True,color=C_DARK,wrap=True)
    # 結論ボックス
    rect(sl,6.28,5.7,6.7,1.18,fill=RGBColor(240,230,255),radius=True)
    rect(sl,6.28,5.7,0.07,1.18,fill=C_PURPLE)
    txt(sl,'💡 結論  →  複勝で少額バック推奨',6.45,5.75,6.35,0.4,size=11,bold=True,color=C_PURPLE)
    txt(sl,'血統最高 × データ消し = 複勝狙い一択',
        6.45,6.2,6.35,0.6,size=10,color=C_DARK)

    pic(sl,IMG_GINYO,9.8,3.8,3.3,2.48)  # ギーニョ右中
    footer(sl,11,paid=True)

# ══ Slide 12: 続きは... (CM) ═══════════════════
def s12(prs):
    sl=add_slide(prs)
    bg(sl,C_PINK)
    pic(sl,IMG_BRAIN,3.5,0.0,10.0,7.5)     # 脳画像背景
    rect(sl,0,0,13.33,7.5,fill=C_PINK)     # ピンクオーバーレイ

    rect(sl,0,0,13.33,0.08,fill=C_GOLD)
    rect(sl,0,7.42,13.33,0.08,fill=C_GOLD)

    # STOP
    rect(sl,3.0,0.52,7.33,0.62,fill=C_WHITE,radius=True)
    txt(sl,'◆  前半ブロック  ここまで  ◆',3.0,0.52,7.33,0.62,
        size=15,bold=True,color=C_PINK,align=PP_ALIGN.CENTER)

    # 続きは
    txt(sl,'続 き は...',0.5,1.38,12.33,2.2,size=80,bold=True,
        color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_S)

    # 有料
    txt(sl,'有料メンバーシップ  限定公開',0.5,3.62,12.33,0.85,
        size=30,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER,font=FONT_S)

    # 内容予告
    previews=[
        '▶ 全馬スコアランキング完全版（19頭）',
        '▶ 対抗馬・注意馬 詳細ファクター解説',
        '▶ 完全買い目提案（単勝〜3連単まで）',
    ]
    for i,p in enumerate(previews):
        rect(sl,2.2,4.6+i*0.56,8.93,0.48,fill=C_PINK_DK,radius=True)
        txt(sl,p,2.4,4.65+i*0.56,8.5,0.36,size=13,color=C_WHITE)

    # CTA
    rect(sl,3.0,6.44,7.33,0.68,fill=C_GOLD,radius=True)
    txt(sl,'▶  チャンネル登録・メンバーシップはこちら',
        3.1,6.44,7.13,0.68,size=14,bold=True,color=C_DARK,align=PP_ALIGN.CENTER)

    # キャラクター2人
    pic(sl,IMG_GINYO, 0.0,4.6,2.7,2.02)   # ギーニョ左下
    pic(sl,IMG_GOTO,10.4,3.8,2.9,2.18)   # ゴトー右下
    txt(sl,'[チャンネルURL / QRコードをここに挿入]',
        0,7.12,13.33,0.27,size=9,italic=True,color=RGBColor(255,220,235),align=PP_ALIGN.CENTER)

# ══ メイン ═════════════════════════════════════
def main():
    prs=new_prs()
    steps=[(s01,'表紙'),(s02,'①2番人気70%'),(s03,'②G1前走100%'),
           (s04,'③上がり3F80%'),(s05,'④10番人気以下0%'),(s06,'⑤G2前走5%'),
           (s07,'血統解析'),(s08,'特別解析'),(s09,'総合スコア1位'),
           (s10,'単勝スコア1位'),(s11,'穴馬スコア1位'),(s12,'続きは...')]
    for fn,name in steps:
        fn(prs); print(f'✅ {name}')
    out=os.path.join(BASE,'桜花賞徹底解析2026_v3.pptx')
    prs.save(out)
    print(f'\n🎉 完成: {out}')

if __name__=='__main__':
    main()
