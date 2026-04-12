#!/usr/bin/env python3
"""桜花賞 解析プレゼンテーション生成スクリプト (9枚: 無料4+CM1+有料4)"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── ディメンション ──
W_IN = 13.33
H_IN = 7.5

# ── カラーパレット（色は3色＋ニュートラル） ──
C_DARK     = RGBColor(26,  8,  20)   # #1a0814  背景Dark
C_DARK2    = RGBColor(45, 11,  30)   # #2d0b1e  背景Medium
C_CRIMSON  = RGBColor(192, 24, 96)   # #c01860  ブランドピンク
C_CRIMSON2 = RGBColor(123, 16, 64)   # #7b1040  ブランドDark
C_GOLD     = RGBColor(252,211, 77)   # #fcd34d  ゴールド（単勝）
C_GOLD_DK  = RGBColor(180, 83,  9)   # #b45309  アンバー
C_GREEN    = RGBColor( 16,185,129)   # #10b981  グリーン（複勝）
C_GREEN_DK = RGBColor(  4,120, 87)   # #047857
C_PURPLE   = RGBColor(109, 40,217)   # #6d28d9  パープル（穴馬）
C_WHITE    = RGBColor(255,255,255)
C_MUTED    = RGBColor(255,200,220)   # muted pink
C_GRAY_LT  = RGBColor(249,242,245)   # #f9f2f5
C_TEXT     = RGBColor( 26, 10, 16)   # #1a0a10

FONT_SERIF = "游明朝"
FONT_SANS  = "游ゴシック"

# ── データロード ──
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
with open(os.path.join(BASE,'docs/data/race-notes/2026-04-12-hanshin-11r.json'),encoding='utf-8') as f:
    DATA = json.load(f)
RACE = DATA['race']
HORSES_RAW = DATA['horses']

# ── スコア計算（entries.html と同ロジック） ──
GP = {'3S':12,'2S':11,'1S':10,'3A':9,'2A':8,'1A':7,'3B':6,'2B':5,'1B':4,
      '3C':3,'2C':2,'1C':1,'D':0,'S':11,'A':8,'B':5,'C':2}
GM = {'zi':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'ext':{'S':['A','S'],'A':['B','A'],'B':['C','B']},
      'kinso':{'S':['S','S'],'A':['A','A'],'B':['B','A'],'C':['D','C'],'D':['D','D']},
      'pace':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'body':{'S':['A','S'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'course':{'S':['A','S'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'style':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']}}

def zi_g(v):
    if not v: return 'B'
    if v>=130: return 'S'
    if v>=115: return 'A'
    if v>=100: return 'B'
    if v>=85:  return 'C'
    return 'D'

def ext_g(f):
    if not f or f=='在厩調整': return 'B'
    if '天栄' in f or 'しがらき' in f: return 'S'
    if 'チャンピオン' in f or '山元' in f or 'キャニオン' in f: return 'A'
    return 'B'

def splt(g, key):
    if not g or g=='—': return [None,None]
    if key in GM and g in GM[key]: return list(GM[key][g])
    return [g,g] if g else [None,None]

def gp(g): return GP.get(g) if g else None

def calc(h):
    zig=zi_g(h.get('sabcZI')); exg=ext_g(h.get('extFacility'))
    zt,zf=splt(zig,'zi'); et,ef=splt(exg,'ext')
    kt,kf=splt(h.get('kinsoGrade'),'kinso'); pt,pf=splt(h.get('paceGrade'),'pace')
    bdt,bdf=splt(h.get('bodyGrade'),'body'); ct,cf=splt(h.get('courseGrade'),'course')
    syt,syf=splt(h.get('styleGrade'),'style')
    srt=h.get('sireGradeTan'); srf=h.get('sireGradeFuku')
    dst=h.get('damSireGradeTan'); dsf=h.get('damSireGradeFuku')
    sr2t=h.get('sireGrade2Tan','B'); sr2f=h.get('sireGrade2Fuku','B')
    ds2t=h.get('damSireGrade2Tan','B'); ds2f=h.get('damSireGrade2Fuku','B')
    blt=h.get('bloodGradeTan') if h.get('bloodGradeTan') and h.get('bloodGradeTan')!='—' else None
    blf=h.get('bloodGradeFuku') if h.get('bloodGradeFuku') and h.get('bloodGradeFuku')!='—' else None
    stt=h.get('stableGradeTan'); stf=h.get('stableGradeFuku')
    brt=h.get('breederGradeTan'); brf=h.get('breederGradeFuku')
    jkt=h.get('jockeyGradeTan'); jkf=h.get('jockeyGradeFuku')
    gtt=h.get('gateGradeTan'); gtf=h.get('gateGradeFuku')
    spt=h.get('specialGradeTan') or h.get('specialGrade')
    spf=h.get('specialGradeFuku') or h.get('specialGrade')
    agt=h.get('ageGradeTan'); agf=h.get('ageGradeFuku')
    rtt=h.get('rotGradeTan'); rtf=h.get('rotGradeFuku')
    s3t=h.get('sf3GradeTan'); s3f=h.get('sf3GradeFuku')
    cdt=h.get('conditionGradeTan'); cdf=h.get('conditionGradeFuku')
    smt=h.get('somaGradeTan'); smf=h.get('somaGradeFuku')
    tgs=[zt,et,srt,dst,sr2t,ds2t,blt,kt,pt,stt,bdt,ct,syt,brt,jkt,gtt,spt,agt,rtt,s3t,cdt,smt]
    fgs=[zf,ef,srf,dsf,sr2f,ds2f,blf,kf,pf,stf,bdf,cf,syf,brf,jkf,gtf,spf,agf,rtf,s3f,cdf,smf]
    ts=tf=tm=fm=0
    for g in tgs:
        p=gp(g)
        if p is not None: ts+=p; tm+=12
    for g in fgs:
        p=gp(g)
        if p is not None: tf+=p; fm+=12
    sf3=(h.get('sf3AnaScore') or 0); rab=(h.get('raceAnaBonus') or 0)
    btg=blt or 'B'; bfg=blf or 'B'
    bp=((gp(btg) or 5)+(gp(bfg) or 5)-10)/24
    ana=sf3+rab+bp
    return {'tanScore':ts,'tanMax':tm,'fukuScore':tf,'fukuMax':fm,'anaScore':ana}

horse_data=[]
for name,h in HORSES_RAW.items():
    sc=calc(h)
    tp=sc['tanScore']/sc['tanMax'] if sc['tanMax']>0 else 0
    fp=sc['fukuScore']/sc['fukuMax'] if sc['fukuMax']>0 else 0
    horse_data.append({**sc,'name':name,'h':h,'tanPct':tp,'fukuPct':fp,'odds':h.get('expectedOdds',999)})

by_tan  = sorted(horse_data,key=lambda x:-x['tanPct'])
by_fuku = sorted(horse_data,key=lambda x:-x['fukuPct'])
by_ana  = sorted(horse_data,key=lambda x:-x['anaScore'])
by_odds = sorted(horse_data,key=lambda x:x['odds'])
for i,d in enumerate(by_tan):  d['tanRank']=i+1
for i,d in enumerate(by_fuku): d['fukuRank']=i+1
for i,d in enumerate(by_ana):  d['anaRank']=i+1
n2d={d['name']:d for d in horse_data}

# レース判定
def detect_type():
    p0=by_tan[0]['tanPct']; p1=by_tan[1]['tanPct']
    p2=by_tan[2]['tanPct']; p3=by_tan[3]['tanPct']
    if (p0-p1)>=0.05: return '1強戦'
    if (p0-p1)<0.05 and (p1-p2)>=0.05: return '2強戦'
    if (p0-p1)<0.05 and (p1-p2)<0.07 and (p2-p3)>=0.05: return '3強戦'
    return '混戦'
RACE_TYPE=detect_type()

# ── PPTX ヘルパー ──
def new_prs():
    prs=Presentation()
    prs.slide_width=Inches(W_IN)
    prs.slide_height=Inches(H_IN)
    return prs

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill=slide.background.fill
    fill.solid()
    fill.fore_color.rgb=color

def rect(slide, x,y,w,h, fill=None, line=None, lw=1.0, radius=False):
    shp_type=5 if radius else 1
    s=slide.shapes.add_shape(shp_type,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb=fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb=line; s.line.width=Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x,y,w,h, size=14, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, font=None, wrap=True):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    run=p.add_run(); run.text=text
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
    if color: run.font.color.rgb=color
    fn=font or FONT_SANS
    run.font.name=fn
    # East Asian font も設定
    rPr=run._r.get_or_add_rPr()
    ea=rPr.find(qn('a:ea'))
    if ea is None: ea=etree.SubElement(rPr,qn('a:ea'))
    ea.set('typeface',fn)
    return tb

def badge_rect(slide, label, x,y, fill=None, text_color=None, size=11, w=None):
    """ラベル付き小バッジ"""
    from pptx.util import Inches as I, Pt as P
    fw=w or (len(label)*0.13+0.25)
    fh=0.32
    fc=fill or C_CRIMSON
    tc=text_color or C_WHITE
    rect(slide, x,y, fw,fh, fill=fc)
    txt(slide, label, x+0.04,y+0.02, fw-0.06,fh-0.02, size=size, bold=True, color=tc, align=PP_ALIGN.CENTER)
    return fw

def horz_bar(slide, x,y,bar_w, pct, fill_color, track_color=None, h=0.12):
    """横棒グラフ"""
    tc=track_color or RGBColor(60,20,40)
    rect(slide, x,y,bar_w,h, fill=tc)
    if pct>0:
        rect(slide, x,y, bar_w*min(pct,1.0), h, fill=fill_color)

# ── Slide 1: 表紙 ──────────────────────────────────────────
def slide1(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    # 左サイドアクセント
    rect(sl, 0,0,0.18,7.5, fill=C_CRIMSON)
    # 上部装飾ライン
    rect(sl, 0.18,0,13.15,0.06, fill=C_CRIMSON2)
    # 右ゴールド縦ライン
    rect(sl, 12.8,0,0.08,7.5, fill=C_GOLD)
    # G1ゴールドバッジ
    rect(sl, 1.4,1.6,1.4,0.55, fill=C_GOLD)
    txt(sl,"G1",1.4,1.6,1.4,0.55, size=22,bold=True,color=C_DARK,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    # メインタイトル
    txt(sl,"桜　花　賞",1.4,2.2,10,1.9, size=82,bold=True,color=C_WHITE,align=PP_ALIGN.LEFT,font=FONT_SERIF)
    # サブタイトル
    txt(sl,"2026年4月12日（日）　阪神競馬場　芝1600m外回り　19頭立て",
        1.4,4.3,10,0.55, size=16,color=C_MUTED,align=PP_ALIGN.LEFT)
    # 仕切り線
    rect(sl, 1.4,5.0,10.5,0.04, fill=C_CRIMSON)
    # キャッチコピー
    txt(sl,"ドリームコア vs スターアニス　—— スコアが示す真の一番人気を解剖する",
        1.4,5.15,10.5,0.55, size=15,italic=True,color=C_GOLD,align=PP_ALIGN.LEFT)
    # ボトムバー
    rect(sl, 0,7.05,13.33,0.45, fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026",
        0.3,7.1,12,0.35, size=12,color=RGBColor(255,180,210),align=PP_ALIGN.LEFT)
    txt(sl,"無料公開",10.8,7.08,2.3,0.38, size=13,bold=True,
        color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 2: レース判定 ─────────────────────────────────────
def slide2(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_CRIMSON)
    # タイトル
    txt(sl,"レース判定",0.35,0.22,5,0.6,size=32,bold=True,color=C_WHITE,font=FONT_SERIF)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_CRIMSON)
    # 判定バッジ
    type_colors={'1強戦':C_GOLD,'2強戦':C_GREEN,'3強戦':C_CRIMSON,'混戦':C_PURPLE}
    tc=type_colors.get(RACE_TYPE,C_CRIMSON)
    rect(sl,0.35,1.05,2.2,0.7,fill=tc)
    txt(sl,RACE_TYPE,0.35,1.05,2.2,0.7,size=26,bold=True,color=C_DARK,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    # 判定コメント
    t1=by_tan[0]; t2=by_tan[1]
    comment=f"単勝スコア1位「{t1['name']}」({t1['tanPct']*100:.1f}%)が2位「{t2['name']}」({t2['tanPct']*100:.1f}%)を{((t1['tanPct']-t2['tanPct'])*100):.1f}pt上回る"
    txt(sl,comment,2.7,1.1,10,0.6,size=14,color=C_MUTED)
    # スコアバーチャート (TOP6)
    txt(sl,"▍ 単勝スコア TOP6",0.35,1.92,5,0.4,size=13,bold=True,color=C_GOLD)
    txt(sl,"複勝スコア TOP6 ▶",7.2,1.92,5.5,0.4,size=13,bold=True,color=C_GREEN,align=PP_ALIGN.RIGHT)
    bar_x=0.35; bar_w=5.8; row_h=0.72
    for i,d in enumerate(by_tan[:6]):
        y=2.42+i*row_h
        # ランクバッジ
        rnk_col=[C_GOLD,RGBColor(156,163,175),RGBColor(194,65,12),C_DARK2,C_DARK2,C_DARK2][i]
        rect(sl,bar_x,y,0.42,0.52,fill=rnk_col)
        txt(sl,f"#{i+1}",bar_x,y,0.42,0.52,size=14,bold=True,color=C_DARK if i<3 else C_MUTED,align=PP_ALIGN.CENTER)
        # 馬名
        txt(sl,d['name'],bar_x+0.5,y,2.3,0.3,size=12,bold=(i==0),color=C_WHITE if i==0 else C_MUTED)
        txt(sl,f"{d['odds']}倍",bar_x+0.5,y+0.28,1.0,0.24,size=9,color=RGBColor(255,200,100))
        # バー
        horz_bar(sl,bar_x+2.9,y+0.12,2.7,d['tanPct'],C_GOLD,h=0.28)
        txt(sl,f"{d['tanPct']*100:.1f}%",bar_x+5.65,y+0.1,0.7,0.28,size=10,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)
    # 複勝 RIGHT side
    bar_rx=7.2; bar_rw=5.5
    for i,d in enumerate(by_fuku[:6]):
        y=2.42+i*row_h
        rnk_col=[C_GREEN,RGBColor(156,163,175),RGBColor(194,65,12),C_DARK2,C_DARK2,C_DARK2][i]
        rect(sl,bar_rx,y,0.42,0.52,fill=rnk_col)
        txt(sl,f"#{i+1}",bar_rx,y,0.42,0.52,size=14,bold=True,color=C_DARK if i<3 else C_MUTED,align=PP_ALIGN.CENTER)
        txt(sl,d['name'],bar_rx+0.5,y,2.3,0.3,size=12,bold=(i==0),color=C_WHITE if i==0 else C_MUTED)
        txt(sl,f"{d['odds']}倍",bar_rx+0.5,y+0.28,1.0,0.24,size=9,color=RGBColor(255,200,100))
        horz_bar(sl,bar_rx+2.9,y+0.12,2.1,d['fukuPct'],C_GREEN,h=0.28)
        txt(sl,f"{d['fukuPct']*100:.1f}%",bar_rx+5.05,y+0.1,0.55,0.28,size=10,bold=True,color=C_GREEN,align=PP_ALIGN.RIGHT)
    # ボトム
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  2/9",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"無料公開",10.8,7.08,2.3,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 3: 注目馬TOP3 ─────────────────────────────────────
def slide3(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_CRIMSON)
    txt(sl,"注目馬ランキング",0.35,0.22,8,0.6,size=32,bold=True,color=C_WHITE,font=FONT_SERIF)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_CRIMSON)
    col_defs=[
        ('単勝ランク','単勝スコア上位',C_GOLD,C_GOLD_DK,by_tan,'tanPct'),
        ('複勝ランク','複勝スコア上位',C_GREEN,C_GREEN_DK,by_fuku,'fukuPct'),
        ('穴馬ランク','穴馬スコア上位',C_PURPLE,RGBColor(76,29,149),by_ana,'anaScore'),
    ]
    col_x=[0.35,4.85,9.35]; col_w=4.3
    for ci,(label,sub,col,col_dk,lst,key) in enumerate(col_defs):
        cx=col_x[ci]
        # カラムヘッダー
        rect(sl,cx,1.08,col_w,0.55,fill=col_dk)
        txt(sl,label,cx,1.08,col_w,0.55,size=18,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_SERIF)
        # TOP3 カード
        card_colors=[col_dk,C_DARK2,C_DARK2]
        rnk_colors=[col,RGBColor(156,163,175),RGBColor(180,90,30)]
        rank_labels=["🥇 1位","🥈 2位","🥉 3位"]
        for ri,hd in enumerate(lst[:3]):
            cy=1.75+ri*1.72
            rect(sl,cx,cy,col_w,1.6,fill=card_colors[ri],radius=True)
            rect(sl,cx,cy,col_w,0.44,fill=rnk_colors[ri])
            txt(sl,rank_labels[ri],cx+0.1,cy+0.04,2.0,0.38,size=12,bold=True,color=C_DARK,align=PP_ALIGN.LEFT)
            txt(sl,f"{hd['odds']}倍",cx+col_w-1.2,cy+0.04,1.1,0.38,size=13,bold=True,color=C_DARK,align=PP_ALIGN.RIGHT)
            txt(sl,hd['name'],cx+0.12,cy+0.52,col_w-0.24,0.48,size=17,bold=True,color=C_WHITE,align=PP_ALIGN.LEFT,font=FONT_SERIF)
            sc_val=hd[key]*100 if key!='anaScore' else hd[key]
            sc_str=f"スコア {sc_val:.1f}{'%' if key!='anaScore' else 'pt'}"
            txt(sl,sc_str,cx+0.12,cy+1.08,col_w-0.24,0.3,size=11,color=col,align=PP_ALIGN.LEFT)
            prev=hd['h'].get('prevRaceName','')
            txt(sl,f"前走：{prev}",cx+0.12,cy+1.35,col_w-0.24,0.25,size=9,italic=True,color=C_MUTED)
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  3/9",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"無料公開",10.8,7.08,2.3,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 4: 全頭スコア一覧 ──────────────────────────────────
def slide4(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_CRIMSON)
    txt(sl,"全頭スコア一覧",0.35,0.22,8,0.6,size=32,bold=True,color=C_WHITE,font=FONT_SERIF)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_CRIMSON)
    # ヘッダー
    heads=[("馬　名",2.8),("単勝%",1.0),("複勝%",1.0),("穴スコ",1.0),("オッズ",1.0),("想定人気",1.25)]
    hx=0.35
    header_y=1.02
    rect(sl,hx,header_y,12.5,0.42,fill=C_CRIMSON2)
    for hd,hw in heads:
        txt(sl,hd,hx+0.05,header_y+0.05,hw-0.1,0.35,size=10,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER)
        hx+=hw
    # 行データ（単勝順）
    row_h=0.42; table_y=1.46
    row_colors=[RGBColor(35,12,25),C_DARK2]
    sorted_for_table=by_tan  # 単勝ランク順
    # 19頭 → 2列 (10+9) で表示
    col_groups=[(sorted_for_table[:10],0.35),(sorted_for_table[10:],6.6)]
    for start_x_in,group in [(0.35,sorted_for_table[:10]),(6.6,sorted_for_table[10:])]:
        for ri,d in enumerate(group):
            ry=table_y+ri*row_h
            row_c=row_colors[ri%2]
            rect(sl,start_x_in,ry,5.9,row_h-0.02,fill=row_c)
            # ランクハイライト
            if d['tanRank']==1:
                rect(sl,start_x_in,ry,0.06,row_h-0.02,fill=C_GOLD)
            elif d['tanRank']<=3:
                rect(sl,start_x_in,ry,0.06,row_h-0.02,fill=C_CRIMSON)
            name_c=C_GOLD if d['tanRank']==1 else (C_WHITE if d['tanRank']<=3 else C_MUTED)
            txt(sl,f"#{d['tanRank']:2d} {d['name']}",start_x_in+0.1,ry+0.06,2.5,0.32,size=10,bold=(d['tanRank']==1),color=name_c)
            cols=[
                (f"{d['tanPct']*100:.1f}%", C_GOLD),
                (f"{d['fukuPct']*100:.1f}%", C_GREEN),
                (f"{d['anaScore']:.2f}", C_PURPLE),
                (f"{d['odds']}倍", C_WHITE),
                (f"{d['h'].get('expectedOdds',999):.1f}→", C_MUTED),
            ]
            cx=start_x_in+2.65
            for val,vc in cols:
                txt(sl,val,cx,ry+0.06,0.88,0.32,size=9,bold=(vc==C_GOLD and d['tanRank']==1),color=vc,align=PP_ALIGN.CENTER)
                cx+=0.9
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  4/9",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"無料公開 ここまで",10.3,7.08,2.8,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 5: コマーシャル ────────────────────────────────────
def slide5(prs):
    sl=add_slide(prs)
    # フルブラック系グラデ代わりの矩形群
    bg(sl,RGBColor(10,2,8))
    rect(sl,0,0,13.33,3.8,fill=RGBColor(26,8,20))
    rect(sl,0,0,0.5,7.5,fill=C_CRIMSON)
    rect(sl,12.83,0,0.5,7.5,fill=C_CRIMSON)
    rect(sl,0,0,13.33,0.08,fill=C_GOLD)
    rect(sl,0,7.42,13.33,0.08,fill=C_GOLD)
    # STOP mark
    rect(sl,4.5,0.55,4.33,0.55,fill=C_CRIMSON)
    txt(sl,"◆  無料公開はここまで  ◆",4.5,0.55,4.33,0.55,size=14,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    # 主テキスト
    txt(sl,"続きは...",0.7,1.35,11.93,1.4,size=72,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    # サブ
    rect(sl,1.2,3.1,10.93,0.08,fill=RGBColor(100,30,60))
    txt(sl,"有料メンバーシップ限定公開",1.2,3.28,10.93,0.75,size=28,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    # コンテンツ予告
    previews=["▶ 主役馬・対抗馬 完全詳細分析","▶ 穴馬3頭の詳細ファクター解説","▶ 完全買い目提案（単勝〜3連単まで）"]
    for i,p in enumerate(previews):
        py=4.18+i*0.52
        rect(sl,2.5,py,8.33,0.44,fill=RGBColor(45,11,30))
        rect(sl,2.5,py,0.06,0.44,fill=C_GOLD)
        txt(sl,p,2.7,py+0.06,7.9,0.32,size=14,color=C_WHITE)
    # CTA
    rect(sl,3.0,6.25,7.33,0.65,fill=C_CRIMSON)
    txt(sl,"▶ チャンネル登録・メンバーシップ登録はこちら",3.0,6.25,7.33,0.65,
        size=15,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,"[チャンネルURL / QRコードをここに挿入]",0.7,7.0,11.93,0.38,
        size=11,italic=True,color=RGBColor(200,150,170),align=PP_ALIGN.CENTER)

# ── Slide 6: 有料① 主役馬詳細 ──────────────────────────────
def slide6(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_GOLD)
    txt(sl,"主役馬 詳細分析",0.35,0.22,8,0.6,size=32,bold=True,color=C_GOLD,font=FONT_SERIF)
    txt(sl,"🔒 メンバーシップ限定",9.0,0.28,4.1,0.45,size=13,bold=True,color=C_CRIMSON,align=PP_ALIGN.RIGHT)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_GOLD)
    # 2頭並列: 単勝1位 vs 単勝2位
    for ci,(hd,cx) in enumerate([(by_tan[0],0.35),(by_tan[1],6.85)]):
        h=hd['h']; cw=6.2
        # カードBG
        rect(sl,cx,1.05,cw,5.85,fill=C_DARK2,radius=True)
        rect(sl,cx,1.05,cw,0.65,fill=C_CRIMSON)
        rank_lbl=f"単勝 #{hd['tanRank']}位　{hd['odds']}倍"
        txt(sl,rank_lbl,cx+0.12,1.1,cw-0.24,0.52,size=12,bold=True,color=C_WHITE)
        txt(sl,hd['name'],cx+0.12,1.78,cw-0.24,0.65,size=24,bold=True,color=C_GOLD,font=FONT_SERIF)
        # スコアバー
        txt(sl,"単勝",cx+0.12,2.52,0.55,0.25,size=9,color=C_GOLD)
        horz_bar(sl,cx+0.7,2.55,3.5,hd['tanPct'],C_GOLD,h=0.18)
        txt(sl,f"{hd['tanPct']*100:.1f}%",cx+4.25,2.52,0.6,0.25,size=9,bold=True,color=C_GOLD)
        txt(sl,"複勝",cx+0.12,2.82,0.55,0.25,size=9,color=C_GREEN)
        horz_bar(sl,cx+0.7,2.85,3.5,hd['fukuPct'],C_GREEN,h=0.18)
        txt(sl,f"{hd['fukuPct']*100:.1f}%",cx+4.25,2.82,0.6,0.25,size=9,bold=True,color=C_GREEN)
        rect(sl,cx+0.12,3.12,cw-0.24,0.03,fill=RGBColor(80,30,55))
        # 特注ポイント
        sp=h.get('specialNote','')
        if sp:
            txt(sl,"⭐ 特注ポイント",cx+0.12,3.2,cw-0.24,0.3,size=10,bold=True,color=C_GOLD)
            txt(sl,sp,cx+0.12,3.52,cw-0.24,0.85,size=9,italic=True,color=C_MUTED,wrap=True)
        # 前走コメント
        pjc=h.get('prevJockeyComment','')[:80] if h.get('prevJockeyComment') else ''
        if pjc:
            txt(sl,"▍ 騎手コメント",cx+0.12,4.42,cw-0.24,0.3,size=10,bold=True,color=C_CRIMSON)
            txt(sl,pjc+'…',cx+0.12,4.72,cw-0.24,0.75,size=8,color=C_MUTED,wrap=True)
        # 調教評価
        cond=h.get('conditionNote','')[:70] if h.get('conditionNote') else ''
        if cond:
            txt(sl,"▍ 調教",cx+0.12,5.52,cw-0.24,0.28,size=10,bold=True,color=RGBColor(139,195,74))
            txt(sl,cond+'…',cx+0.12,5.8,cw-0.24,0.65,size=8,color=C_MUTED,wrap=True)
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  6/9  🔒",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"有料限定",10.8,7.08,2.3,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 7: 有料② 対抗馬詳細 ──────────────────────────────
def slide7(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_GREEN)
    txt(sl,"対抗馬・注意馬 詳細分析",0.35,0.22,9,0.6,size=30,bold=True,color=C_GREEN,font=FONT_SERIF)
    txt(sl,"🔒 メンバーシップ限定",9.0,0.28,4.1,0.45,size=13,bold=True,color=C_CRIMSON,align=PP_ALIGN.RIGHT)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_GREEN)
    # 3〜5位の馬を3列で表示
    trio=[by_tan[2],by_tan[3],by_tan[4]]
    col_x=[0.35,4.62,8.89]; col_w=4.05
    for ci,hd in enumerate(trio):
        cx=col_x[ci]; h=hd['h']
        rect(sl,cx,1.05,col_w,5.85,fill=C_DARK2,radius=True)
        rect(sl,cx,1.05,col_w,0.55,fill=C_DARK2)
        rect(sl,cx,1.05,col_w,0.55,fill=RGBColor(4,80,55))
        txt(sl,f"単勝#{hd['tanRank']}  {hd['odds']}倍",cx+0.1,1.1,col_w-0.2,0.42,size=11,bold=True,color=C_WHITE)
        txt(sl,hd['name'],cx+0.1,1.68,col_w-0.2,0.58,size=20,bold=True,color=C_GREEN,font=FONT_SERIF)
        # バー
        txt(sl,"単",cx+0.1,2.35,0.28,0.22,size=8,color=C_GOLD)
        horz_bar(sl,cx+0.42,2.37,3.15,hd['tanPct'],C_GOLD,h=0.16)
        txt(sl,f"{hd['tanPct']*100:.1f}%",cx+3.6,2.35,0.42,0.22,size=8,color=C_GOLD,align=PP_ALIGN.RIGHT)
        txt(sl,"複",cx+0.1,2.62,0.28,0.22,size=8,color=C_GREEN)
        horz_bar(sl,cx+0.42,2.64,3.15,hd['fukuPct'],C_GREEN,h=0.16)
        txt(sl,f"{hd['fukuPct']*100:.1f}%",cx+3.6,2.62,0.42,0.22,size=8,color=C_GREEN,align=PP_ALIGN.RIGHT)
        rect(sl,cx+0.1,2.95,col_w-0.2,0.03,fill=RGBColor(30,70,50))
        # 特注
        sp=h.get('specialNote','')
        if sp:
            txt(sl,"⭐ "+sp,cx+0.1,3.06,col_w-0.2,0.7,size=8,italic=True,color=C_GOLD,wrap=True)
        # コメント
        pjc=(h.get('prevJockeyComment') or '')[:100]
        if pjc:
            txt(sl,"▍ "+pjc+'…',cx+0.1,3.82,col_w-0.2,1.5,size=7,color=C_MUTED,wrap=True)
        # 血統コメント
        bn=(h.get('bloodNote') or '')[:60]
        if bn:
            txt(sl,"🧬 "+bn,cx+0.1,5.38,col_w-0.2,0.65,size=7,italic=True,color=RGBColor(196,181,253),wrap=True)
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  7/9  🔒",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"有料限定",10.8,7.08,2.3,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 8: 有料③ 穴馬詳細 ──────────────────────────────
def slide8(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_PURPLE)
    txt(sl,"穴馬 詳細分析",0.35,0.22,8,0.6,size=32,bold=True,color=RGBColor(196,181,253),font=FONT_SERIF)
    txt(sl,"🔒 メンバーシップ限定",9.0,0.28,4.1,0.45,size=13,bold=True,color=C_CRIMSON,align=PP_ALIGN.RIGHT)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_PURPLE)
    # 穴馬TOP3
    col_x=[0.35,4.62,8.89]; col_w=4.05
    for ci,hd in enumerate(by_ana[:3]):
        cx=col_x[ci]; h=hd['h']
        rect(sl,cx,1.05,col_w,5.85,fill=C_DARK2,radius=True)
        rect(sl,cx,1.05,col_w,0.55,fill=RGBColor(60,20,120))
        txt(sl,f"穴#{ci+1}  {hd['odds']}倍  (単#{hd['tanRank']})",cx+0.1,1.1,col_w-0.2,0.42,size=11,bold=True,color=RGBColor(196,181,253))
        txt(sl,hd['name'],cx+0.1,1.68,col_w-0.2,0.58,size=20,bold=True,color=RGBColor(196,181,253),font=FONT_SERIF)
        txt(sl,f"穴スコア: {hd['anaScore']:.2f}pt",cx+0.1,2.34,col_w-0.2,0.32,size=11,color=C_PURPLE)
        # 穴バー
        ana_max=by_ana[0]['anaScore']
        horz_bar(sl,cx+0.1,2.72,col_w-0.2,min(hd['anaScore']/max(ana_max,0.01),1.0),C_PURPLE,h=0.2)
        rect(sl,cx+0.1,3.0,col_w-0.2,0.03,fill=RGBColor(60,20,80))
        # 穴馬ファクター
        ran=h.get('raceAnaNote','')
        if ran:
            txt(sl,"🎯 "+ran,cx+0.1,3.1,col_w-0.2,0.45,size=9,bold=True,color=C_GOLD,wrap=True)
        sp=h.get('specialNote','')
        if sp:
            txt(sl,"⭐ "+sp,cx+0.1,3.62,col_w-0.2,0.7,size=8,italic=True,color=C_GOLD,wrap=True)
        pjc=(h.get('prevJockeyComment') or '')[:90]
        if pjc:
            txt(sl,"▍ "+pjc+'…',cx+0.1,4.4,col_w-0.2,1.2,size=7,color=C_MUTED,wrap=True)
        bn=(h.get('bloodNote') or '')[:60]
        if bn:
            txt(sl,"🧬 "+bn,cx+0.1,5.65,col_w-0.2,0.55,size=7,italic=True,color=RGBColor(196,181,253),wrap=True)
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  8/9  🔒",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"有料限定",10.8,7.08,2.3,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── Slide 9: 有料④ 買い目提案 ──────────────────────────────
def slide9(prs):
    sl=add_slide(prs); bg(sl,C_DARK)
    rect(sl,0,0,0.18,7.5,fill=C_GOLD)
    txt(sl,"買い目提案",0.35,0.22,7,0.6,size=32,bold=True,color=C_GOLD,font=FONT_SERIF)
    txt(sl,"🔒 メンバーシップ限定",9.0,0.28,4.1,0.45,size=13,bold=True,color=C_CRIMSON,align=PP_ALIGN.RIGHT)
    rect(sl,0.35,0.9,12.5,0.04,fill=C_GOLD)
    # レース判定
    rect(sl,0.35,1.0,3.5,0.58,fill=C_GOLD)
    txt(sl,RACE_TYPE,0.35,1.0,3.5,0.58,size=24,bold=True,color=C_DARK,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    t0=by_tan[0]; t1=by_tan[1]; f0=by_fuku[0]; a0=by_ana[0]
    txt(sl,f"単勝1位：{t0['name']} / 複勝1位：{f0['name']} / 穴馬1位：{a0['name']}",
        4.0,1.1,9.0,0.42,size=12,color=C_MUTED)
    rect(sl,0.35,1.66,12.5,0.04,fill=RGBColor(80,30,55))
    # 買い目テーブル
    bets=[
        ("必 須","単 勝","1点", f"{t0['name']}",C_GOLD,"単勝ランク1位 固定"),
        ("必 須","複 勝","2点", f"{f0['name']} / {a0['name']}",C_GREEN,"複勝1位＋穴馬1位 固定（重複時は繰り下げ）"),
        ("2連系","馬 連","6点", f"軸：{t0['name']}  ×  相手：上位6頭",C_GOLD,"1強戦 → 軸1頭流し6点"),
        ("3連系","3連単","30点",f"1着固定：{t0['name']}  ×  2,3着：上位6頭",C_CRIMSON,"1強戦 → 1着固定30点フォーメーション"),
    ]
    row_y=1.82; row_h=1.2; rw=12.5
    for bi,(cat,bet_type,pts,horses,bc,note) in enumerate(bets):
        ry=row_y+bi*row_h
        row_col=RGBColor(35,12,25) if bi%2==0 else C_DARK2
        rect(sl,0.35,ry,rw,row_h-0.05,fill=row_col)
        rect(sl,0.35,ry,0.06,row_h-0.05,fill=bc)
        # カテゴリバッジ
        rect(sl,0.45,ry+0.12,0.75,0.36,fill=bc)
        txt(sl,cat,0.45,ry+0.12,0.75,0.36,size=9,bold=True,color=C_DARK,align=PP_ALIGN.CENTER)
        # 券種
        rect(sl,1.28,ry+0.12,0.95,0.36,fill=RGBColor(60,20,45))
        txt(sl,bet_type,1.28,ry+0.12,0.95,0.36,size=11,bold=True,color=bc,align=PP_ALIGN.CENTER)
        # 点数
        txt(sl,pts,2.32,ry+0.12,0.8,0.36,size=16,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
        # 馬名
        txt(sl,horses,3.2,ry+0.1,8.5,0.42,size=12,bold=True,color=C_WHITE)
        # 備考
        txt(sl,note,3.2,ry+0.54,9.2,0.38,size=9,italic=True,color=C_MUTED)
    # 合計
    rect(sl,0.35,6.72,12.5,0.28,fill=C_CRIMSON2)
    txt(sl,"合計 39点  ／  ※枠番・騎手スコア確定後に微調整あり",
        0.5,6.73,12.2,0.25,size=11,bold=True,color=C_GOLD)
    rect(sl,0,7.05,13.33,0.45,fill=C_CRIMSON2)
    txt(sl,"UG競馬 競馬予想チャンネル  ／  桜花賞解析レポート 2026  ———  9/9  🔒",
        0.3,7.1,12.7,0.35,size=10,color=RGBColor(255,180,210))
    txt(sl,"有料限定",10.8,7.08,2.3,0.38,size=12,bold=True,color=C_GOLD,align=PP_ALIGN.RIGHT)

# ── メイン ──────────────────────────────────────────────────
def main():
    prs=new_prs()
    print(f"レース判定: {RACE_TYPE}")
    print(f"単勝1位: {by_tan[0]['name']} ({by_tan[0]['tanPct']*100:.1f}%)")
    print(f"複勝1位: {by_fuku[0]['name']} ({by_fuku[0]['fukuPct']*100:.1f}%)")
    print(f"穴馬1位: {by_ana[0]['name']} ({by_ana[0]['anaScore']:.2f}pt)")
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)
    out=os.path.join(os.path.dirname(BASE),'桜花賞解析2026.pptx') if False else os.path.join(BASE,'桜花賞解析2026.pptx')
    prs.writeFile=None  # pptx uses save()
    prs.save(out)
    print(f"\n✅ 保存完了: {out}")

if __name__=='__main__':
    main()
