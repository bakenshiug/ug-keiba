#!/usr/bin/env python3
"""桜花賞徹底解析 プレゼンテーション v2
カラー: ビビッドピンク / ホワイト / クリーム
構成: 前半8枚（無料4+有料解説3+CM1）
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── ディメンション ──
W_IN, H_IN = 13.33, 7.5

# ── カラーパレット（ビビッドピンク / ホワイト / クリーム） ──
C_PINK    = RGBColor(232, 19, 110)   # #E8136E  ビビッドピンク（支配色）
C_PINK_DK = RGBColor(180,  8, 80)    # #B40850  ダークピンク
C_PINK_LT = RGBColor(255,100,175)    # #FF64AF  ライトピンク
C_CREAM   = RGBColor(255,251,245)    # #FFFBF5  クリーム白（コンテンツBG）
C_WHITE   = RGBColor(255,255,255)    # #FFFFFF
C_GOLD    = RGBColor(245,183, 49)    # #F5B731  ゴールド（脳画像から）
C_GOLD_DK = RGBColor(180, 83,  9)    # #B45309
C_DARK    = RGBColor( 30, 10, 20)    # #1E0A14  テキスト色
C_GRAY    = RGBColor(120, 80,100)    # #785064  サブテキスト
C_GREEN   = RGBColor( 16,185,129)    # #10B981  複勝グリーン
C_PURPLE  = RGBColor(109, 40,217)    # #6D28D9  穴馬パープル

FONT_SERIF = "游明朝"
FONT_SANS  = "游ゴシック"

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_BRAIN = os.path.join(BASE, 'brain_cover.jpg')

# ── データ ──
with open(os.path.join(BASE,'docs/data/race-notes/2026-04-12-hanshin-11r.json'),encoding='utf-8') as f:
    DATA = json.load(f)
RACE_ANA = DATA.get('raceAnalysis','')
BLOOD_ANA = DATA.get('bloodAnalysis','')
HORSES_RAW = DATA['horses']

# ── スコア計算 ──
GP = {'3S':12,'2S':11,'1S':10,'3A':9,'2A':8,'1A':7,'3B':6,'2B':5,'1B':4,
      '3C':3,'2C':2,'1C':1,'D':0,'S':11,'A':8,'B':5,'C':2}
GM = {'zi':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'ext':{'S':['A','S'],'A':['B','A'],'B':['C','B']},
      'kinso':{'S':['S','S'],'A':['A','A'],'B':['B','A'],'C':['D','C'],'D':['D','D']},
      'pace':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'body':{'S':['A','S'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'course':{'S':['A','S'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']},
      'style':{'S':['S','A'],'A':['A','A'],'B':['C','B'],'C':['D','C'],'D':['D','D']}}

def zi_g(v):
    if not v: return 'B'
    return 'S' if v>=130 else 'A' if v>=115 else 'B' if v>=100 else 'C' if v>=85 else 'D'

def ext_g(f):
    if not f or f=='在厩調整': return 'B'
    if '天栄' in f or 'しがらき' in f: return 'S'
    if 'チャンピオン' in f or '山元' in f or 'キャニオン' in f: return 'A'
    return 'B'

def splt(g, key):
    if not g or g=='—': return [None,None]
    if key in GM and g in GM[key]: return list(GM[key][g])
    return [g,g] if g else [None,None]

def gp(g): return GP.get(g) if g else None

def calc(h):
    zig=zi_g(h.get('sabcZI')); exg=ext_g(h.get('extFacility'))
    zt,zf=splt(zig,'zi'); et,ef=splt(exg,'ext')
    kt,kf=splt(h.get('kinsoGrade'),'kinso'); pt,pf=splt(h.get('paceGrade'),'pace')
    bdt,bdf=splt(h.get('bodyGrade'),'body'); ct,cf=splt(h.get('courseGrade'),'course')
    syt,syf=splt(h.get('styleGrade'),'style')
    srt=h.get('sireGradeTan'); srf=h.get('sireGradeFuku')
    dst=h.get('damSireGradeTan'); dsf=h.get('damSireGradeFuku')
    sr2t=h.get('sireGrade2Tan','B'); sr2f=h.get('sireGrade2Fuku','B')
    ds2t=h.get('damSireGrade2Tan','B'); ds2f=h.get('damSireGrade2Fuku','B')
    blt=h.get('bloodGradeTan') if h.get('bloodGradeTan') and h.get('bloodGradeTan')!='—' else None
    blf=h.get('bloodGradeFuku') if h.get('bloodGradeFuku') and h.get('bloodGradeFuku')!='—' else None
    stt=h.get('stableGradeTan'); stf=h.get('stableGradeFuku')
    brt=h.get('breederGradeTan'); brf=h.get('breederGradeFuku')
    jkt=h.get('jockeyGradeTan'); jkf=h.get('jockeyGradeFuku')
    gtt=h.get('gateGradeTan'); gtf=h.get('gateGradeFuku')
    spt=h.get('specialGradeTan') or h.get('specialGrade')
    spf=h.get('specialGradeFuku') or h.get('specialGrade')
    agt=h.get('ageGradeTan'); agf=h.get('ageGradeFuku')
    rtt=h.get('rotGradeTan'); rtf=h.get('rotGradeFuku')
    s3t=h.get('sf3GradeTan'); s3f=h.get('sf3GradeFuku')
    cdt=h.get('conditionGradeTan'); cdf=h.get('conditionGradeFuku')
    smt=h.get('somaGradeTan'); smf=h.get('somaGradeFuku')
    tgs=[zt,et,srt,dst,sr2t,ds2t,blt,kt,pt,stt,bdt,ct,syt,brt,jkt,gtt,spt,agt,rtt,s3t,cdt,smt]
    fgs=[zf,ef,srf,dsf,sr2f,ds2f,blf,kf,pf,stf,bdf,cf,syf,brf,jkf,gtf,spf,agf,rtf,s3f,cdf,smf]
    ts=tf=tm=fm=0
    for g in tgs:
        p=gp(g)
        if p is not None: ts+=p; tm+=12
    for g in fgs:
        p=gp(g)
        if p is not None: tf+=p; fm+=12
    sf3=(h.get('sf3AnaScore') or 0); rab=(h.get('raceAnaBonus') or 0)
    btg=blt or 'B'; bfg=blf or 'B'
    bp=((gp(btg) or 5)+(gp(bfg) or 5)-10)/24
    ana=sf3+rab+bp
    sougou=(ts+tf)/(tm+fm) if (tm+fm)>0 else 0
    return {'tanScore':ts,'tanMax':tm,'fukuScore':tf,'fukuMax':fm,'anaScore':ana,'sougouScore':sougou}

horse_data=[]
for name,h in HORSES_RAW.items():
    sc=calc(h)
    tp=sc['tanScore']/sc['tanMax'] if sc['tanMax']>0 else 0
    fp=sc['fukuScore']/sc['fukuMax'] if sc['fukuMax']>0 else 0
    horse_data.append({**sc,'name':name,'h':h,'tanPct':tp,'fukuPct':fp,'odds':h.get('expectedOdds',999)})

by_tan    = sorted(horse_data,key=lambda x:-x['tanPct'])
by_fuku   = sorted(horse_data,key=lambda x:-x['fukuPct'])
by_ana    = sorted(horse_data,key=lambda x:-x['anaScore'])
by_sougou = sorted(horse_data,key=lambda x:-x['sougouScore'])
for i,d in enumerate(by_tan):    d['tanRank']=i+1
for i,d in enumerate(by_fuku):   d['fukuRank']=i+1
for i,d in enumerate(by_ana):    d['anaRank']=i+1
for i,d in enumerate(by_sougou): d['sougouRank']=i+1

# ── PPTX ヘルパー ──
def new_prs():
    prs=Presentation(); prs.slide_width=Inches(W_IN); prs.slide_height=Inches(H_IN)
    return prs

def add_slide(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill=slide.background.fill; fill.solid(); fill.fore_color.rgb=color

def rect(slide, x,y,w,h, fill=None, line=None, lw=1.0, radius=False):
    s=slide.shapes.add_shape(5 if radius else 1,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    if line: s.line.color.rgb=line; s.line.width=Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, x,y,w,h, size=14, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, font=None, wrap=True):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    run=p.add_run(); run.text=text
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
    if color: run.font.color.rgb=color
    fn=font or FONT_SANS; run.font.name=fn
    rPr=run._r.get_or_add_rPr()
    ea=rPr.find(qn('a:ea'))
    if ea is None: ea=etree.SubElement(rPr,qn('a:ea'))
    ea.set('typeface',fn)
    return tb

def img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def pink_header(slide, title_text, subtitle=None):
    """ピンクヘッダー帯 + タイトル"""
    rect(slide, 0,0,13.33,1.45, fill=C_PINK)
    txt(slide, title_text, 0.5,0.18,10,0.92, size=34,bold=True,color=C_WHITE,font=FONT_SERIF)
    if subtitle:
        txt(slide, subtitle, 0.5,0.98,12,0.44, size=12,italic=True,color=RGBColor(255,210,230))
    # ページ番号スロット（呼び出し元で上書き）

def footer(slide, page_num, total=8, paid=False):
    rect(slide, 0,7.1,13.33,0.4, fill=C_PINK_DK)
    tag="🔒 有料限定" if paid else "🆓 無料公開"
    tag_col=C_GOLD if not paid else RGBColor(255,220,130)
    txt(slide,"UG競馬 競馬予想チャンネル  ／  桜花賞徹底解析 2026",
        0.4,7.13,10,0.3,size=10,color=RGBColor(255,210,230))
    txt(slide,f"{tag}  {page_num}/{total}",9.5,7.13,3.5,0.3,size=10,bold=True,
        color=tag_col,align=PP_ALIGN.RIGHT)

def stat_card(slide, x,y,w,h, number, label, sub=None, bg_col=None, num_col=None):
    """大きな数字カード"""
    bc=bg_col or C_WHITE
    nc=num_col or C_PINK
    rect(slide, x,y,w,h, fill=bc, radius=True)
    rect(slide, x,y,w,0.04, fill=nc)
    txt(slide, number, x+0.1,y+0.1,w-0.2,h*0.55, size=28,bold=True,color=nc,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    txt(slide, label, x+0.05,y+h*0.58,w-0.1,h*0.28, size=10,bold=True,color=C_DARK,align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, x+0.05,y+h*0.82,w-0.1,h*0.2, size=8,italic=True,color=C_GRAY,align=PP_ALIGN.CENTER)

def horz_bar(slide, x,y,bar_w,pct,fill_color,track_color=None,h=0.14):
    tc=track_color or RGBColor(230,210,220)
    rect(slide, x,y,bar_w,h, fill=tc)
    if pct>0:
        rect(slide, x,y, bar_w*min(pct,1.0),h, fill=fill_color)

# ══════════════════════════════════════════════
# Slide 1: 表紙
# ══════════════════════════════════════════════
def slide1(prs):
    sl=add_slide(prs)
    bg(sl, C_PINK)

    # 右半分に脳画像（右端にはみ出すイメージ）
    img(sl, IMG_BRAIN, 5.8, -0.1, 7.8, 7.7)

    # 左側オーバーレイ（白からピンクへのフェード代わりの帯）
    rect(sl, 0,0, 6.5,7.5, fill=C_PINK)

    # 上部ゴールドライン
    rect(sl, 0,0, 13.33,0.07, fill=C_GOLD)
    # 下部ゴールドライン
    rect(sl, 0,7.43,13.33,0.07, fill=C_GOLD)

    # G1バッジ
    rect(sl, 0.45,0.65, 1.2,0.5, fill=C_GOLD)
    txt(sl,"G  1",0.45,0.65,1.2,0.5,size=18,bold=True,color=C_DARK,align=PP_ALIGN.CENTER,font=FONT_SERIF)

    # メインタイトル
    txt(sl,"桜花賞",0.4,1.3,6.0,2.0,size=72,bold=True,color=C_WHITE,align=PP_ALIGN.LEFT,font=FONT_SERIF)

    # サブタイトル
    rect(sl, 0.4,3.42,5.6,0.06, fill=C_GOLD)
    txt(sl,"徹底解析",0.4,3.55,5.6,0.9,size=38,bold=True,color=C_GOLD,align=PP_ALIGN.LEFT,font=FONT_SERIF)

    # 日時・場所
    txt(sl,"2026年4月12日（日）  阪神競馬場  芝1600m外  19頭立て",
        0.4,4.62,5.8,0.48,size=13,color=RGBColor(255,220,235))

    # 区切り線
    rect(sl, 0.4,5.22,5.6,0.05, fill=RGBColor(255,180,210))

    # キャッチ
    txt(sl,"ギーニョ重賞データ×血統×特別解析",
        0.4,5.35,5.8,0.5,size=14,italic=True,color=RGBColor(255,240,248))

    # クリーム帯ボトム
    rect(sl, 0,6.88,6.6,0.6, fill=RGBColor(255,245,252))
    txt(sl,"UG競馬 競馬予想チャンネル",0.45,6.93,5.8,0.45,size=12,color=C_PINK_DK,bold=True)

# ══════════════════════════════════════════════
# Slide 2: レースのポイント
# ══════════════════════════════════════════════
def slide2(prs):
    sl=add_slide(prs)
    bg(sl, C_CREAM)
    pink_header(sl,"レースのポイント","ギーニョ 重賞データ解析より　過去10年 桜花賞 徹底分析")

    # 5つのデータカード (Z字配置)
    cards=[
        ("70%","2番人気の複勝率","1番人気(60%)を上回る最強人気"),
        ("100%","G1前走1着の複勝率","阪神JF直行ローテ 過去10年全複勝圏"),
        ("80%","上がり3F 1位の複勝率","直線で切れる馬が主役のレース"),
        ("0%","10番人気以下の複勝率","人気薄は完全消し。上位人気に集中"),
        ("5%","G2前走1着の複勝率","チューリップ賞前走は危険ゾーン"),
    ]
    card_w=2.35; card_h=2.15; gap=0.18
    start_x=0.42; start_y=1.65
    colors=[C_PINK, C_GREEN, C_GOLD_DK, C_GRAY, RGBColor(220,38,38)]
    for i,(num,lbl,sub) in enumerate(cards):
        cx=start_x+i*(card_w+gap)
        stat_card(sl,cx,start_y,card_w,card_h,num,lbl,sub,
                  bg_col=C_WHITE,num_col=colors[i])
        # 下部カラーバー
        rect(sl,cx,start_y+card_h-0.06,card_w,0.06,fill=colors[i])

    # 追加インサイト
    rect(sl, 0.42,3.98,12.5,0.05, fill=RGBColor(240,220,228))
    insights=[
        "✅ 中9週以上が複勝率34.3%と優位　— 阪神JF直行は最高条件",
        "⚠️ アネモネS/フラワーS/紅梅S前走は複勝率0%の完全消し",
        "✅ エルフィンS前走は複勝率40%　— フィリーズレビューは6.8%",
        "✅ 馬体重460〜479kgがベストゾーン（複勝29.4%）",
    ]
    for i,ins in enumerate(insights):
        col=C_PINK if '✅' in ins else RGBColor(200,50,50)
        rect(sl, 0.42,4.1+i*0.68,0.04,0.44,fill=col)
        txt(sl,ins,0.58,4.12+i*0.68,12.2,0.46,size=12,color=C_DARK,bold=('✅' in ins and i==0))

    footer(sl,2)

# ══════════════════════════════════════════════
# Slide 3: 血統
# ══════════════════════════════════════════════
def slide3(prs):
    sl=add_slide(prs)
    bg(sl, C_CREAM)
    pink_header(sl,"血　統","ギーニョ 血統データ解析　父系・母父系 複勝率ランキング")

    # 左：父系
    rect(sl, 0.4,1.6,5.9,0.5,fill=C_PINK)
    txt(sl,"父　系　ランキング",0.4,1.6,5.9,0.5,size=16,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_SERIF)

    sires=[
        ("キタサンブラック","66.7%","3頭","ブラックチャリス該当"),
        ("エピファネイア",   "28.6%","7頭","アランカール該当"),
        ("ロードカナロア",   "22.2%","勝率0%","複は取れるが勝てない血統"),
    ]
    for i,(sire,pct,cnt,note) in enumerate(sires):
        ry=2.2+i*1.22
        pct_val=float(pct.replace('%',''))
        is_top=(i==0)
        rc=C_WHITE
        rect(sl,0.4,ry,5.9,1.08,fill=rc,radius=True)
        rect(sl,0.4,ry,0.07,1.08,fill=C_PINK if not is_top else C_GOLD)
        rank_c=[C_GOLD,C_PINK,C_GRAY][i]
        txt(sl,f"#{i+1}",0.55,ry+0.08,0.5,0.5,size=16,bold=True,color=rank_c,align=PP_ALIGN.CENTER)
        txt(sl,sire,1.12,ry+0.1,2.9,0.42,size=14,bold=is_top,color=C_DARK,font=FONT_SERIF)
        txt(sl,f"複勝率 {pct}  ({cnt})",1.12,ry+0.55,3.2,0.3,size=10,color=C_GRAY)
        horz_bar(sl,4.4,ry+0.25,1.65,pct_val/100,C_PINK if not is_top else C_GOLD,h=0.2)
        txt(sl,note,1.12,ry+0.8,4.7,0.22,size=8,italic=True,color=C_PINK)

    # 右：母父系
    rect(sl, 6.8,1.6,6.0,0.5,fill=C_PINK_DK)
    txt(sl,"母父系　ランキング",6.8,1.6,6.0,0.5,size=16,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_SERIF)

    damsires=[
        ("ハーツクライ",  "50.0%","2頭","フェスティバルヒル該当"),
        ("クロフネ",       "42.9%","7頭",""),
        ("ディープインパクト","0.0%","多数","ルールザウェイヴ要注意"),
    ]
    for i,(ds,pct,cnt,note) in enumerate(damsires):
        ry=2.2+i*1.22
        pct_val=float(pct.replace('%',''))
        is_top=(i==0)
        rect(sl,6.8,ry,6.0,1.08,fill=C_WHITE,radius=True)
        rect(sl,6.8,ry,0.07,1.08,fill=C_GOLD if is_top else C_PINK)
        rank_c=[C_GOLD,C_PINK,C_GRAY][i]
        txt(sl,f"#{i+1}",6.95,ry+0.08,0.5,0.5,size=16,bold=True,color=rank_c,align=PP_ALIGN.CENTER)
        txt(sl,ds,7.52,ry+0.1,3.2,0.42,size=14,bold=is_top,color=C_DARK,font=FONT_SERIF)
        txt(sl,f"複勝率 {pct}  ({cnt})",7.52,ry+0.55,3.4,0.3,size=10,color=C_GRAY)
        horz_bar(sl,10.9,ry+0.25,1.62,pct_val/100,C_GOLD if is_top else C_PINK,h=0.2)
        if note:
            txt(sl,note,7.52,ry+0.8,5.1,0.22,size=8,italic=True,color=C_PINK)

    # 血統総括
    rect(sl, 0.4,5.88,12.5,0.06,fill=C_PINK)
    rect(sl, 0.4,6.0,12.5,0.98,fill=C_WHITE)
    rect(sl, 0.4,6.0,0.07,0.98,fill=C_GOLD)
    txt(sl,"血統総括",0.6,6.04,2,0.3,size=11,bold=True,color=C_PINK)
    txt(sl,"血統トップ：ブラックチャリス（父キタサンブラック66.7%）とフェスティバルヒル（母父ハーツクライ50%）。"
        "ただしブラックチャリスはレースデータとの「データの相克」あり。取捨は慎重に。",
        0.6,6.35,12.2,0.58,size=11,color=C_DARK,wrap=True)

    footer(sl,3)

# ══════════════════════════════════════════════
# Slide 4: 特別解析
# ══════════════════════════════════════════════
def slide4(prs):
    sl=add_slide(prs)
    bg(sl, C_CREAM)
    pink_header(sl,"特別解析","ギーニョ重賞データ×specialNote 特別ファクター該当馬一覧")

    # 特注馬ピックアップ（specialNoteがある馬）
    special_horses=[(name,h) for name,h in HORSES_RAW.items() if h.get('specialNote')]

    # 上段: 特注馬カード2列
    n2d_map={d['name']:d for d in horse_data}
    col_x=[0.4,6.86]; col_w=6.06
    for ci, (name,h) in enumerate(special_horses[:2]):
        cx=col_x[ci]; hd=n2d_map.get(name,{})
        odds=h.get('expectedOdds',999)
        rect(sl,cx,1.6,col_w,2.05,fill=C_WHITE,radius=True)
        rect(sl,cx,1.6,col_w,0.5,fill=C_PINK)
        txt(sl,f"⭐ 特注  {odds}倍",cx+0.1,1.65,col_w-0.2,0.4,size=13,bold=True,color=C_WHITE)
        txt(sl,name,cx+0.1,2.18,col_w-0.2,0.55,size=20,bold=True,color=C_PINK,font=FONT_SERIF)
        sp=h.get('specialNote','')
        txt(sl,sp,cx+0.1,2.8,col_w-0.2,0.72,size=9,italic=True,color=C_DARK,wrap=True)

    # 中段: G2前走消し解説
    rect(sl,0.4,3.82,12.5,0.06,fill=RGBColor(220,38,38))
    rect(sl,0.4,3.95,12.5,1.1,fill=C_WHITE)
    rect(sl,0.4,3.95,0.07,1.1,fill=RGBColor(220,38,38))
    txt(sl,"⚠️ G2前走消し vs G1前走推し（データからの結論）",
        0.6,3.99,12,0.36,size=13,bold=True,color=RGBColor(180,20,20))
    txt(sl,"G2前走1着（チューリップ賞など）の複勝率わずか5%。どんなに実力があっても割り引くのがデータの正解。"
        "一方でG1前走（阪神JF等）は複勝率53.8%、前走1着は100%。前走クラスは最重要ファクター。",
        0.6,4.42,12.2,0.55,size=11,color=C_DARK,wrap=True)

    # 消し馬リスト
    rect(sl,0.4,5.15,12.5,0.38,fill=RGBColor(255,240,240))
    keshi_horses=[n for n,h in HORSES_RAW.items() if h.get('prevRaceName','') in
                  ['アネモネS','フラワーS','紅梅S']]
    keshi_str="前走消し（複勝率0%ローテ）: "+", ".join(keshi_horses) if keshi_horses else "前走消しローテ該当馬 要確認"
    txt(sl,keshi_str,0.6,5.22,12,0.28,size=10,color=RGBColor(180,20,20))

    # raceAnaNote 該当馬
    rect(sl,0.4,5.62,12.5,0.4,fill=C_PINK)
    txt(sl,"穴馬ファクター該当馬（raceAnaBonus 1.0以上）",0.6,5.67,12,0.3,size=12,bold=True,color=C_WHITE)
    ana_horses=[(name,h.get('raceAnaNote',''),h.get('raceAnaBonus',0))
                for name,h in HORSES_RAW.items() if (h.get('raceAnaBonus') or 0)>=1.0]
    ana_str="  /  ".join(f"{n}（{note}）" for n,note,_ in sorted(ana_horses,key=lambda x:-x[2]))
    rect(sl,0.4,6.06,12.5,0.95,fill=C_WHITE)
    txt(sl,ana_str,0.55,6.1,12.3,0.85,size=10,color=C_DARK,wrap=True)

    footer(sl,4)

# ══════════════════════════════════════════════
# Slide 5: 総合スコア1位
# ══════════════════════════════════════════════
def slide5(prs):
    sl=add_slide(prs)
    bg(sl, C_CREAM)
    top=by_sougou[0]; h=top['h']
    pink_header(sl,f"総合スコア 1位",f"単勝×複勝 バランス総合評価  —  {top['name']}  {top['odds']}倍")

    # 左: 大きな馬名カード
    rect(sl, 0.4,1.6,5.5,5.35,fill=C_WHITE)
    rect(sl, 0.4,1.6,5.5,0.55,fill=C_PINK)
    txt(sl,"総合スコア  #1",0.55,1.65,5.2,0.45,size=14,bold=True,color=C_WHITE)
    txt(sl,top['name'],0.5,2.25,5.3,1.0,size=36,bold=True,color=C_PINK,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    txt(sl,f"想定オッズ {top['odds']}倍",0.5,3.32,5.3,0.38,size=14,color=C_GRAY,align=PP_ALIGN.CENTER)

    # スコアバー3本
    for i,(label,pct,col) in enumerate([
        ("総合",top['sougouScore'],C_PINK),
        ("単勝",top['tanPct'],C_GOLD_DK),
        ("複勝",top['fukuPct'],C_GREEN)]):
        by=3.88+i*0.62
        txt(sl,label,0.55,by,0.6,0.42,size=10,bold=True,color=col)
        horz_bar(sl,1.22,by+0.1,3.5,pct,col,h=0.22)
        txt(sl,f"{pct*100:.1f}%",4.78,by,0.85,0.42,size=11,bold=True,color=col,align=PP_ALIGN.RIGHT)

    # 特注ポイント
    sp=h.get('specialNote','')
    if sp:
        rect(sl, 0.55,5.78,5.2,0.06,fill=C_GOLD)
        txt(sl,"⭐  "+sp,0.55,5.9,5.2,0.88,size=9,italic=True,color=C_DARK,wrap=True)

    # 右: ファクター詳細
    rect(sl, 6.2,1.6,6.7,5.35,fill=C_WHITE)
    rect(sl, 6.2,1.6,6.7,0.55,fill=C_PINK_DK)
    txt(sl,"ファクター詳細",6.38,1.65,6.3,0.45,size=14,bold=True,color=C_WHITE)

    factors=[
        ("能力(ZI)",     f"{h.get('sabcZI','?')}  (最上位)" if h.get('sabcZI','?')>=130 else str(h.get('sabcZI','?'))),
        ("近走",         str(h.get('kinsoGrade','—'))),
        ("コース適性",   str(h.get('courseGrade','—'))),
        ("ローテ",       str(h.get('rotGradeTan','—'))),
        ("厩舎",         str(h.get('stableGradeTan','—'))),
        ("調教",         str(h.get('conditionGradeTan','—'))),
        ("前走",         str(h.get('prevRaceName','—'))),
        ("前走結果",     str(h.get('prevFinish','—'))),
    ]
    for i,(label,val) in enumerate(factors):
        fy=2.28+i*0.56
        bc=RGBColor(255,240,246) if i%2==0 else C_WHITE
        rect(sl,6.2,fy,6.7,0.52,fill=bc)
        txt(sl,label,6.35,fy+0.08,2.4,0.36,size=10,color=C_GRAY)
        is_good= val in ['S','2S','3S','A','2A','3A'] or '1着' in val or '最上位' in val
        vc=C_PINK if is_good else C_DARK
        txt(sl,val,8.85,fy+0.08,3.9,0.36,size=11,bold=is_good,color=vc,align=PP_ALIGN.RIGHT)

    # 条件コメント
    cond=h.get('conditionNote','')
    if cond:
        rect(sl, 6.2,6.03,6.7,0.85,fill=RGBColor(255,245,252))
        rect(sl, 6.2,6.03,0.06,0.85,fill=C_GOLD)
        txt(sl,cond[:110]+'…',6.35,6.07,6.45,0.72,size=8,italic=True,color=C_DARK,wrap=True)

    footer(sl,5,paid=True)

# ══════════════════════════════════════════════
# Slide 6: 単勝スコア1位
# ══════════════════════════════════════════════
def slide6(prs):
    sl=add_slide(prs)
    bg(sl, C_CREAM)
    top=by_tan[0]; h=top['h']

    # ゴールドヘッダー（単勝=ゴールド）
    rect(sl,0,0,13.33,1.45,fill=C_GOLD_DK)
    txt(sl,"単勝スコア 1位",0.5,0.18,10,0.92,size=34,bold=True,color=C_WHITE,font=FONT_SERIF)
    txt(sl,f"単勝確率最上位馬  —  {top['name']}  {top['odds']}倍",
        0.5,0.98,12,0.44,size=12,italic=True,color=RGBColor(255,240,200))

    # 上段2カラム
    # 左: 単勝スコアビジュアル
    rect(sl, 0.4,1.6,5.5,2.4,fill=C_WHITE)
    txt(sl,top['name'],0.5,1.7,5.3,1.0,size=30,bold=True,color=C_GOLD_DK,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    txt(sl,f"単勝スコア  {top['tanPct']*100:.1f}%",0.5,2.75,5.3,0.45,size=16,bold=True,color=C_GOLD_DK,align=PP_ALIGN.CENTER)
    horz_bar(sl,0.6,3.28,5.1,top['tanPct'],C_GOLD_DK,h=0.26)
    txt(sl,f"19頭中 単勝1位 / 複勝{top['fukuRank']}位",0.5,3.62,5.3,0.32,size=10,color=C_GRAY,align=PP_ALIGN.CENTER)

    # 右: なぜ単勝1位か
    rect(sl, 6.2,1.6,6.7,2.4,fill=C_WHITE)
    rect(sl, 6.2,1.6,6.7,0.42,fill=C_GOLD_DK)
    txt(sl,"なぜ単勝1位か？",6.35,1.65,6.3,0.35,size=13,bold=True,color=C_WHITE)
    reasons=[]
    if h.get('sabcZI',0)>=130: reasons.append(f"ZI={h['sabcZI']}（全馬最上位）")
    if (h.get('rotGradeTan') or '') in ['S','2S','3S']: reasons.append(f"ローテ{h.get('rotGradeTan')}（最高ローテ）")
    if (h.get('conditionGradeTan') or '') in ['S','2S','3S']: reasons.append(f"調教{h.get('conditionGradeTan')}（万全仕上がり）")
    if (h.get('courseGrade') or '') in ['S','A']: reasons.append(f"コース適性{h.get('courseGrade')}")
    if h.get('specialNote'): reasons.append("特注ファクター該当")
    for ri,r in enumerate(reasons[:4]):
        rect(sl,6.35,2.12+ri*0.47,6.4,0.4,fill=RGBColor(255,248,224) if ri%2==0 else C_WHITE)
        txt(sl,f"✓ {r}",6.5,2.16+ri*0.47,6.1,0.32,size=11,color=C_GOLD_DK)

    # 下段: 騎手コメント + 特注
    rect(sl, 0.4,4.15,12.5,0.5,fill=C_PINK)
    txt(sl,"騎手コメント（前走）",0.55,4.2,12,0.38,size=12,bold=True,color=C_WHITE)
    pjc=h.get('prevJockeyComment','')
    rect(sl, 0.4,4.7,12.5,1.25,fill=C_WHITE)
    rect(sl, 0.4,4.7,0.06,1.25,fill=C_GOLD_DK)
    txt(sl,pjc[:200]+'…' if len(pjc)>200 else pjc,
        0.58,4.76,12.2,1.12,size=10,color=C_DARK,italic=True,wrap=True)

    # 特注サマリー
    sp=h.get('specialNote','')
    if sp:
        rect(sl, 0.4,6.04,12.5,0.46,fill=RGBColor(255,250,230))
        rect(sl, 0.4,6.04,0.06,0.46,fill=C_GOLD)
        txt(sl,"⭐ "+sp,0.58,6.09,12.2,0.35,size=10,italic=True,color=C_DARK,wrap=True)

    footer(sl,6,paid=True)

# ══════════════════════════════════════════════
# Slide 7: 穴馬スコア1位
# ══════════════════════════════════════════════
def slide7(prs):
    sl=add_slide(prs)
    bg(sl, C_CREAM)
    top=by_ana[0]; h=top['h']

    # パープルヘッダー
    rect(sl,0,0,13.33,1.45,fill=C_PURPLE)
    txt(sl,"穴馬スコア 1位",0.5,0.18,10,0.92,size=34,bold=True,color=C_WHITE,font=FONT_SERIF)
    txt(sl,f"穴馬スコア最上位  —  {top['name']}  {top['odds']}倍  (単{top['tanRank']}位)",
        0.5,0.98,12,0.44,size=12,italic=True,color=RGBColor(220,200,255))

    # 2カラム構成
    # 左: プロフィール + スコア
    rect(sl, 0.4,1.6,5.5,5.35,fill=C_WHITE)
    rect(sl, 0.4,1.6,5.5,0.5,fill=C_PURPLE)
    txt(sl,"穴馬スコア  #1",0.55,1.65,5.2,0.4,size=13,bold=True,color=C_WHITE)
    txt(sl,top['name'],0.5,2.18,5.3,0.88,size=28,bold=True,color=C_PURPLE,align=PP_ALIGN.CENTER,font=FONT_SERIF)
    txt(sl,f"{top['odds']}倍  (想定{int(round(top['odds']))}番人気前後)",0.5,3.1,5.3,0.38,size=12,color=C_GRAY,align=PP_ALIGN.CENTER)

    for i,(label,pct,col) in enumerate([
        ("穴スコア",min(top['anaScore']/3.0,1.0),C_PURPLE),
        ("単勝スコア",top['tanPct'],C_GOLD_DK),
        ("複勝スコア",top['fukuPct'],C_GREEN)]):
        by=3.62+i*0.6
        txt(sl,label,0.55,by,1.5,0.42,size=9,bold=True,color=col)
        horz_bar(sl,2.15,by+0.1,3.1,pct,col,h=0.2)
        val=f"{top['anaScore']:.2f}pt" if i==0 else f"{pct*100:.1f}%"
        txt(sl,val,5.3,by,0.55,0.42,size=10,bold=True,color=col,align=PP_ALIGN.RIGHT)

    # 穴馬ファクター
    ran=h.get('raceAnaNote','')
    if ran:
        rect(sl, 0.55,5.5,5.2,0.06,fill=C_GOLD)
        txt(sl,"🎯 "+ran,0.55,5.62,5.2,0.42,size=10,bold=True,color=C_GOLD_DK)
    sp=h.get('specialNote','')
    if sp:
        txt(sl,"⭐ "+sp,0.55,6.1,5.2,0.72,size=8,italic=True,color=C_DARK,wrap=True)

    # 右: データの相克（血統最高 vs レースデータ消し）
    rect(sl, 6.2,1.6,6.7,5.35,fill=C_WHITE)
    rect(sl, 6.2,1.6,6.7,0.5,fill=C_GOLD_DK)
    txt(sl,"データの相克　—— なぜ穴馬か",6.38,1.65,6.3,0.4,size=12,bold=True,color=C_WHITE)

    # 血統評価
    rect(sl, 6.2,2.18,6.7,1.3,fill=RGBColor(255,248,220))
    rect(sl, 6.2,2.18,0.06,1.3,fill=C_GOLD)
    txt(sl,"血統評価：最高",6.35,2.22,6.3,0.36,size=12,bold=True,color=C_GOLD_DK)
    txt(sl,"父キタサンブラック → 桜花賞複勝率66.7%\n全産駒が複勝圏という驚異的データ",
        6.35,2.62,6.3,0.78,size=10,color=C_DARK,wrap=True)

    # レースデータ消し
    rect(sl, 6.2,3.58,6.7,1.3,fill=RGBColor(255,240,240))
    rect(sl, 6.2,3.58,0.06,1.3,fill=RGBColor(220,38,38))
    txt(sl,"レースデータ：消し根拠あり",6.35,3.62,6.3,0.36,size=12,bold=True,color=RGBColor(180,20,20))
    ran2=h.get('prevRaceName','')
    txt(sl,f"前走：{ran2}\n穴馬ファクター加点があるが、他のデータ消し根拠と交差",
        6.35,4.02,6.3,0.78,size=10,color=C_DARK,wrap=True)

    # 結論
    rect(sl, 6.2,4.98,6.7,1.92,fill=RGBColor(245,230,255))
    rect(sl, 6.2,4.98,0.06,1.92,fill=C_PURPLE)
    txt(sl,"⚡ 結論：少額で複勝押さえ",6.35,5.02,6.3,0.36,size=12,bold=True,color=C_PURPLE)
    bn=h.get('bloodNote','')
    txt(sl,(bn or '血統単体では最上位評価。ただし複数のデータ消し根拠と交差するため単勝は過信禁物。複勝での少額バックが合理的。')[:150],
        6.35,5.44,6.3,1.35,size=9,color=C_DARK,italic=True,wrap=True)

    footer(sl,7,paid=True)

# ══════════════════════════════════════════════
# Slide 8: 続きは...
# ══════════════════════════════════════════════
def slide8(prs):
    sl=add_slide(prs)
    bg(sl, C_PINK)

    # 脳画像を薄く全面配置（装飾）
    img(sl, IMG_BRAIN, 3.0, 0.0, 10.5, 7.5)

    # 半透明ピンクオーバーレイ
    rect(sl, 0,0,13.33,7.5,fill=C_PINK)

    # 上下ゴールドライン
    rect(sl, 0,0,13.33,0.08,fill=C_GOLD)
    rect(sl, 0,7.42,13.33,0.08,fill=C_GOLD)

    # STOP バナー
    rect(sl, 2.8,0.55,7.73,0.6,fill=C_WHITE)
    txt(sl,"◆  前半ブロック  ここまで  ◆",2.8,0.55,7.73,0.6,
        size=16,bold=True,color=C_PINK,align=PP_ALIGN.CENTER)

    # メインテキスト
    txt(sl,"続きは...",1.0,1.35,11.33,2.1,size=80,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font=FONT_SERIF)

    # サブテキスト
    txt(sl,"有料メンバーシップ　限定公開",1.0,3.55,11.33,0.85,
        size=30,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER,font=FONT_SERIF)

    # 予告コンテンツ
    previews=[
        "▶ 全頭スコアランキング完全版（19頭）",
        "▶ 対抗馬・注意馬 詳細ファクター解説",
        "▶ 完全買い目提案（単勝〜3連単まで）",
    ]
    for i,p in enumerate(previews):
        py=4.62+i*0.58
        rect(sl,2.2,py,8.93,0.5,fill=RGBColor(180,8,80))
        txt(sl,p,2.4,py+0.07,8.5,0.35,size=13,color=C_WHITE)

    # CTAボタン
    rect(sl,3.0,6.5,7.33,0.72,fill=C_GOLD)
    txt(sl,"▶ チャンネル登録・メンバーシップ登録はこちら",
        3.0,6.5,7.33,0.72,size=14,bold=True,color=C_DARK,align=PP_ALIGN.CENTER)

    txt(sl,"[チャンネルURL・QRコードをここに挿入]",
        0,7.12,13.33,0.28,size=10,italic=True,color=RGBColor(255,220,235),align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# メイン
# ══════════════════════════════════════════════
def main():
    prs=new_prs()
    print("生成開始...")
    slide1(prs); print("✅ Slide1: 表紙")
    slide2(prs); print("✅ Slide2: レースのポイント")
    slide3(prs); print("✅ Slide3: 血統")
    slide4(prs); print("✅ Slide4: 特別解析")
    slide5(prs); print("✅ Slide5: 総合スコア1位")
    slide6(prs); print("✅ Slide6: 単勝スコア1位")
    slide7(prs); print("✅ Slide7: 穴馬スコア1位")
    slide8(prs); print("✅ Slide8: 続きは...")
    out=os.path.join(BASE,'桜花賞徹底解析2026.pptx')
    prs.save(out)
    print(f"\n✅ 完成: {out}")

if __name__=='__main__':
    main()
